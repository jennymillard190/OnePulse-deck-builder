import os
from typing import List, Tuple, Dict, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from . import config
from .scale_helpers import calculate_net_group_percentages, format_net_group_line
import logging
import re

# Configure logging
logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)
def clean_chart_title(title: str) -> str:
    """
    Remove prefixes like 'Q(4)' and 'Question:' from chart titles.
    """
    title = str(title).strip()
    title = re.sub(r'^Question:\s*', '', title, flags=re.IGNORECASE)
    title = re.sub(r'^Q\(\d+\)\s*', '', title, flags=re.IGNORECASE)
    return title.strip()

def apply_chart_title(chart, title: str) -> None:
    clean_title = clean_chart_title(title)
    chart.chart_title.text_frame.text = clean_title
    p = chart.chart_title.text_frame.paragraphs[0]
    p.font.name = 'Bally Thrill'
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0,0,0)
    for run in p.runs:
        run.font.name = 'Bally Thrill'
        run.font.size = Pt(12)
        run.font.italic = True
        run.font.color.rgb = RGBColor(0, 0, 0)

def add_source_footer(slide, prs, text: str) -> None:
    """
    Move source slightly up and right so it sits in the white space.
    """
    y = prs.slide_height - Inches(config.FOOTER_OFFSET) - Inches(0.55)
    tx = slide.shapes.add_textbox(Inches(0.75), y, Inches(8), Inches(0.3))
    pf = tx.text_frame.paragraphs[0]
    pf.text = text
    pf.font.name = 'Calibri'
    pf.font.size = Pt(10)
    pf.font.color.rgb = RGBColor(0,0,0)

def add_net_score_callouts(
    slide,
    prs: Presentation,
    categories: List[str],
    series_list: List[Tuple[str, List[float]]]
) -> None:
    """Add compact grouped percentage callouts for recognised scale questions."""
    callouts = []
    for label, values in series_list:
        net_groups = calculate_net_group_percentages(categories, values)
        if net_groups is None:
            return
        callouts.append((label, net_groups))

    if not callouts:
        return

    box_count = len(callouts)
    box_left = prs.slide_width - Inches(2.0)
    box_top = Inches(1.25)
    box_width = Inches(1.6)
    max_total_height = Inches(4.55)
    box_gap = Inches(0.06)
    box_height = min(
        Inches(0.72),
        (max_total_height - (box_count - 1) * box_gap) / box_count
    )
    font_size = Pt(7 if box_height >= Inches(0.53) else 6)

    for idx, (label, net_groups) in enumerate(callouts):
        top = box_top + idx * (box_height + box_gap)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            box_left,
            top,
            box_width,
            box_height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(245, 245, 245)
        box.line.color.rgb = RGBColor(0, 0, 0)
        box.line.width = Pt(0.75)

        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)

        label_para = tf.paragraphs[0]
        label_para.text = str(label)
        label_para.alignment = PP_ALIGN.CENTER
        label_para.font.name = "Bally Thrill"
        label_para.font.size = font_size
        label_para.font.bold = True
        label_para.font.color.rgb = RGBColor(0, 0, 0)

        positive_para = tf.add_paragraph()
        positive_para.text = format_net_group_line(
            str(net_groups["positive_label"]),
            int(net_groups["positive_pct"])
        )
        positive_para.alignment = PP_ALIGN.CENTER
        positive_para.font.name = "Bally Thrill"
        positive_para.font.size = font_size
        positive_para.font.color.rgb = RGBColor(0, 0, 0)

        negative_para = tf.add_paragraph()
        negative_para.text = format_net_group_line(
            str(net_groups["negative_label"]),
            int(net_groups["negative_pct"])
        )
        negative_para.alignment = PP_ALIGN.CENTER
        negative_para.font.name = "Bally Thrill"
        negative_para.font.size = font_size
        negative_para.font.color.rgb = RGBColor(0, 0, 0)

def set_text_style(shape, font_name: str = "Bally Thrill", font_size: int = None, bold: bool = None, italic: bool = None):
    """Apply consistent font styling to all text in a shape."""
    if not shape.has_text_frame:
        return

    for paragraph in shape.text_frame.paragraphs:
        paragraph.font.name = font_name

        if font_size is not None:
            paragraph.font.size = Pt(font_size)
        if bold is not None:
            paragraph.font.bold = bold
        if italic is not None:
            paragraph.font.italic = italic

        paragraph.font.color.rgb = RGBColor(0, 0, 0)

        for run in paragraph.runs:
            run.font.name = font_name

            if font_size is not None:
                run.font.size = Pt(font_size)
            if bold is not None:
                run.font.bold = bold
            if italic is not None:
                run.font.italic = italic

            run.font.color.rgb = RGBColor(0, 0, 0)

def get_layout(prs: Presentation, name: str = "8_UPFRONT 3"):
    """Get a slide layout by name, falling back to a default if not found."""
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    return prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[5]

def clear_text(slide):
    """Remove all text boxes from a slide."""
    # We need to convert to list because we're modifying the collection while iterating
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            shapes_to_remove.append(shape)
    
    # Now remove the shapes
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

def create_chart_slide(
    prs: Presentation,
    cats: List[str],
    series_list: List[Tuple[str, List[float]]]
) -> Tuple[Presentation, 'Chart']:
    """
    Create a clustered column chart with multiple series.
    
    Args:
        prs: Presentation object
        cats: List of category names
        series_list: List of (label, values) tuples, first expected to be Total
    
    Returns:
        Tuple containing:
        - Updated presentation
        - Created chart object
    """
    logger.debug(f"Creating chart with {len(cats)} categories and {len(series_list)} series")
    
    slide = prs.slides.add_slide(get_layout(prs))
    clear_text(slide)
    x = int((prs.slide_width - Inches(config.CHART_WIDTH)) / 2)
    y = int((prs.slide_height - Inches(config.CHART_HEIGHT)) / 2)

    # Create chart data
    cd = ChartData()
    cd.categories = list(cats)
    for label, vals in series_list:
        cd.add_series(label, list(vals))

    # Add chart to slide
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        x, y, Inches(config.CHART_WIDTH), Inches(config.CHART_HEIGHT),
        cd
    ).chart

    # Calculate total number of bars
    total_bars = len(cats) * len(series_list)
    
    # Determine data label font size based on number of bars
    if total_bars > 40:
        data_label_size = Pt(6)
    elif total_bars > 15:
        data_label_size = Pt(10)
    else:
        data_label_size = Pt(12)

    # Define available theme colors in order of preference
    accent_colors = [
        MSO_THEME_COLOR.ACCENT_1,
        MSO_THEME_COLOR.ACCENT_2,
        MSO_THEME_COLOR.ACCENT_3,
        MSO_THEME_COLOR.ACCENT_4,
        MSO_THEME_COLOR.ACCENT_5,
        MSO_THEME_COLOR.ACCENT_6
    ]
    
    secondary_colors = [
        MSO_THEME_COLOR.DARK_1,
        MSO_THEME_COLOR.DARK_2,
        MSO_THEME_COLOR.LIGHT_1,
        MSO_THEME_COLOR.LIGHT_2,
        MSO_THEME_COLOR.BACKGROUND_1,
        MSO_THEME_COLOR.BACKGROUND_2
    ]
    
    # Track used colors to ensure consistent cycling
    used_colors = set()
    
    # Style each series with consistent colors
    for idx, (label, _) in enumerate(series_list):
        ser = chart.series[idx]
        ser.format.fill.solid()
        
        # Assign or retrieve theme color
        if label not in config.theme_lookup:
            # Use a consistent color scheme based on the label
            if label == "Total":
                color = MSO_THEME_COLOR.ACCENT_1  # ACCENT_1 for Total
            else:
                # First try to use accent colors
                accent_index = (hash(label) % len(accent_colors))
                if accent_colors[accent_index] not in used_colors:
                    color = accent_colors[accent_index]
                else:
                    # If accent color is used, find next available accent color
                    for i in range(len(accent_colors)):
                        if accent_colors[i] not in used_colors:
                            color = accent_colors[i]
                            break
                    else:
                        # If all accent colors are used, use secondary colors
                        secondary_index = (hash(label) % len(secondary_colors))
                        if secondary_colors[secondary_index] not in used_colors:
                            color = secondary_colors[secondary_index]
                        else:
                            # If secondary color is used, find next available secondary color
                            for i in range(len(secondary_colors)):
                                if secondary_colors[i] not in used_colors:
                                    color = secondary_colors[i]
                                    break
                            else:
                                # If all colors are used, start reusing accent colors
                                color = accent_colors[0]
                
                used_colors.add(color)
            config.theme_lookup[label] = color
            
        ser.format.fill.fore_color.theme_color = config.theme_lookup[label]
        
        # Format data labels
        dl = ser.data_labels
        dl.show_value = True
        dl.number_format = '0%'
        dl.font.name = 'Bally Thrill'
        dl.font.size = data_label_size

    # Remove gridlines & format axes
    chart.value_axis.visible = False
    chart.value_axis.has_major_gridlines = False
    chart.category_axis.has_major_gridlines = False
    
    # Format category axis labels
    tl = chart.category_axis.tick_labels
    tl.font.name = 'Bally Thrill'
    tl.font.size = Pt(12)
    tl.font.color.rgb = RGBColor(0,0,0)
    tl.number_format = '@'  # Use text format
    tl.orientation = 45  # Angle the labels
    tl.offset = 100  # Add some space between labels and axis
    tl.wrap_text = True  # Enable text wrapping

    # Configure legend
    if len(series_list) > 1:
        chart.has_legend = True
        lg = chart.legend
        lg.position = XL_LEGEND_POSITION.TOP
        lg.include_in_layout = False
        lg.font.name = 'Bally Thrill'
        lg.font.size = Pt(12)
    else:
        chart.has_legend = False

    add_net_score_callouts(slide, prs, cats, series_list)

    return slide, chart

def add_raw_audience_slides(
    prs: Presentation,
    raw_audience_data: List[Tuple[str, List[str], List[Tuple[str, List[float], int]]]]
) -> Presentation:
    """Add raw audience slides to the presentation."""
    for title, categories, segments in raw_audience_data:
        for label, vals, n_resp in segments:
            series = [(label, vals)]
            slide, chart = create_chart_slide(prs, categories, series)
            
            apply_chart_title(chart, title)
            add_source_footer(slide, prs, f"Source: OnePulse, {label} ({n_resp})")
            
    return prs

def add_combined_slides_full_export(
    prs: Presentation,
    combined_data: List[Tuple[str, List[str], List[Tuple[str, List[float], int]]]],
    group_audience_names: set = None
) -> Presentation:
    """Add combined slides for full export: all segments + groups + individual segments."""
    group_audience_names = group_audience_names or set()

    for title, categories, segments in combined_data:
        series_list = []
        footer_parts = ["Source: OnePulse"]

        for label, vals, n_resp in segments:
            series_list.append((label, vals))
            footer_parts.append(f"{label} ({n_resp})")

        # Main combined slide
        slide, chart = create_chart_slide(prs, categories, series_list)
        apply_chart_title(chart, title)
        add_source_footer(slide, prs, ', '.join(footer_parts))

        # Group slides first (if any)
        group_segments = [s for s in segments if s[0] != "Total" and " - " in title]
        if group_segments:
            for label, vals, n_resp in group_segments:
                total_segment = next((s for s in segments if s[0] == "Total"), None)
                if total_segment:
                    segment_series = [
                        (total_segment[0], total_segment[1]),
                        (label, vals)
                    ]
                    slide, chart = create_chart_slide(prs, categories, segment_series)
                    apply_chart_title(chart, title)
                    add_source_footer(
                        slide,
                        prs,
                        f"Source: OnePulse, {label} ({n_resp}), Total ({total_segment[2]})"
                    )

        # Individual segment slides (for ungrouped audiences)
        non_total_segments = [s for s in segments if s[0] != "Total"]
        is_group_chart = " - " in title
        is_individual_chart = title.endswith(")")

        if len(non_total_segments) > 1 and not is_group_chart and not is_individual_chart:
            total_segment = next((s for s in segments if s[0] == "Total"), None)
            if total_segment:
                for label, vals, n_resp in segments:
                    if label != "Total" and label not in group_audience_names:
                        segment_series = [
                            (total_segment[0], total_segment[1]),
                            (label, vals)
                        ]
                        slide, chart = create_chart_slide(prs, categories, segment_series)
                        apply_chart_title(chart, title)
                        add_source_footer(
                            slide,
                            prs,
                            f"Source: OnePulse, {label} ({n_resp}), Total ({total_segment[2]})"
                        )

    return prs


def add_combined_slides_condensed_export(
    prs: Presentation,
    combined_data: List[Tuple[str, List[str], List[Tuple[str, List[float], int]]]],
    group_audience_names: set = None,
    audience_defs: dict = None,
    raw_audience_data: List[Tuple[str, List[str], List[Tuple[str, List[float], int]]]] = None
) -> Presentation:
    """Add combined slides for condensed export: groups + ungrouped only (no duplication)."""
    group_audience_names = group_audience_names or set()
    ungrouped_audiences = set()

    if audience_defs:
        grouped_audiences = set()
        for group in audience_defs.get("__groups__", []):
            grouped_audiences.update(group.get("audiences", []))
        ungrouped_audiences = set(audience_defs.keys()) - grouped_audiences - {"__groups__"}

    # If no audiences are defined, show Total charts only
    if (not audience_defs or (len(audience_defs) == 1 and '__groups__' in audience_defs)) and raw_audience_data:
        for title, categories, segments in raw_audience_data:
            if segments and segments[0][0] == "Total":
                total_segment = segments[0]
                segment_series = [(total_segment[0], total_segment[1])]
                slide, chart = create_chart_slide(prs, categories, segment_series)
                apply_chart_title(chart, title)
                add_source_footer(slide, prs, f"Source: OnePulse, Total ({total_segment[2]})")
        return prs

    slide_count = 0

    for title, categories, segments in combined_data:
        is_group_chart = " - " in title
        is_individual_chart = re.search(r"\([^)]+\)$", title) and not is_group_chart
        non_total_segments = [s for s in segments if s[0] != "Total"]

        if is_group_chart:
            group_name = title.split(" - ")[-1].strip()
            group_members = []

            if audience_defs:
                for group in audience_defs.get("__groups__", []):
                    if group["name"] == group_name:
                        group_members = group.get("audiences", [])
                        break

            total_segment = next((s for s in segments if s[0] == "Total"), None)
            group_segments = [s for s in segments if s[0] != "Total" and s[0] in group_members]

            if total_segment and group_segments:
                segment_series = [(total_segment[0], total_segment[1])]
                group_labels = []
                total_group_resp = 0

                for label, vals, n_resp in group_segments:
                    segment_series.append((label, vals))
                    group_labels.append(label)
                    total_group_resp += n_resp

                combined_group_label = " & ".join(group_labels)
                slide, chart = create_chart_slide(prs, categories, segment_series)
                apply_chart_title(chart, title)
                add_source_footer(
                    slide,
                    prs,
                    f"Source: OnePulse, {combined_group_label} ({total_group_resp}), Total ({total_segment[2]})"
                )
                slide_count += 1

            continue

        # Skip "all segments" charts
        if not is_group_chart and not is_individual_chart and len(non_total_segments) > 1:
            continue

        # Process individual charts for ungrouped audiences
        if is_individual_chart:
            total_segment = next((s for s in segments if s[0] == "Total"), None)
            if total_segment and len(segments) == 2:
                audience_segment = segments[1]
                label, vals, n_resp = audience_segment

                if not audience_defs or label in ungrouped_audiences:
                    segment_series = [
                        (total_segment[0], total_segment[1]),
                        (label, vals)
                    ]
                    slide, chart = create_chart_slide(prs, categories, segment_series)
                    apply_chart_title(chart, title)
                    add_source_footer(
                        slide,
                        prs,
                        f"Source: OnePulse, {label} ({n_resp}), Total ({total_segment[2]})"
                    )
                    slide_count += 1
            continue

        total_segment = next((s for s in segments if s[0] == "Total"), None)
        if total_segment:
            ungrouped_segments = [
                s for s in segments
                if s[0] != "Total" and (not audience_defs or s[0] in ungrouped_audiences)
            ]

            for label, vals, n_resp in ungrouped_segments:
                segment_series = [
                    (total_segment[0], total_segment[1]),
                    (label, vals)
                ]
                slide, chart = create_chart_slide(prs, categories, segment_series)
                apply_chart_title(chart, title)
                add_source_footer(
                    slide,
                    prs,
                    f"Source: OnePulse, {label} ({n_resp}), Total ({total_segment[2]})"
                )
                slide_count += 1

    return prs
    
    slide_count = 0
    for title, categories, segments in combined_data:
        is_group_chart = " - " in title
        is_individual_chart = re.search(r"\([^)]+\)$", title) and not is_group_chart
        non_total_segments = [s for s in segments if s[0] != "Total"]
        if is_group_chart:
            group_name = title.split(" - ")[-1].strip()
            group_members = []
            if audience_defs:
                for group in audience_defs.get("__groups__", []):
                    if group["name"] == group_name:
                        group_members = group.get("audiences", [])
                        break
            segment_labels = [s[0] for s in segments]
            found_members = [m for m in group_members if m in segment_labels]
            total_segment = next((s for s in segments if s[0] == "Total"), None)
            group_segments = [s for s in segments if s[0] != "Total" and s[0] in group_members]
            if total_segment and group_segments:
                segment_series = [(total_segment[0], total_segment[1])]
                group_labels = []
                total_group_resp = 0
                for label, vals, n_resp in group_segments:
                    segment_series.append((label, vals))
                    group_labels.append(label)
                    total_group_resp += n_resp
                combined_group_label = " & ".join(group_labels)
                slide, chart = create_chart_slide(prs, categories, segment_series)
                apply_chart_title(chart, title)
                add_source_footer(slide, prs, f"Source: OnePulse, {label} ({n_resp}), Total ({total_segment[2]})")
                pf.font.name = 'Calibri'
                pf.font.size = Pt(10)
                pf.font.color.rgb = RGBColor(0,0,0)
                slide_count += 1
            continue  # Don't process this chart further
        # Skip 'all segments' charts (not group, not individual, more than one non-Total segment)
        if not is_group_chart and not is_individual_chart and len(non_total_segments) > 1:
            continue
        # Process individual charts for ungrouped audiences
        if is_individual_chart:
            total_segment = next((s for s in segments if s[0] == "Total"), None)
            if total_segment and len(segments) == 2:  # Total + one audience
                audience_segment = segments[1]  # The non-Total segment
                label, vals, n_resp = audience_segment
                # Only process if this audience is ungrouped
                if not audience_defs or label in ungrouped_audiences:
                    segment_series = [
                        (total_segment[0], total_segment[1]),
                        (label, vals)
                    ]
                    slide, chart = create_chart_slide(prs, categories, segment_series)
                apply_chart_title(chart, title)
                add_source_footer(slide, prs, f"Source: OnePulse, {label} ({n_resp}), Total ({total_segment[2]})")
                pf.font.name = 'Calibri'
                pf.font.size = Pt(10)
                pf.font.color.rgb = RGBColor(0,0,0)
                slide_count += 1
            continue
        total_segment = next((s for s in segments if s[0] == "Total"), None)
        if total_segment:
            ungrouped_segments = [
                s for s in segments
                if s[0] != "Total" and (not audience_defs or s[0] in ungrouped_audiences)
            ]
            for label, vals, n_resp in ungrouped_segments:
                segment_series = [
                    (total_segment[0], total_segment[1]),
                    (label, vals)
                ]
                slide, chart = create_chart_slide(prs, categories, segment_series)
                apply_chart_title(chart, title)
                add_source_footer(
                    slide,
                    prs,
                    f"Source: OnePulse, {label} ({n_resp}), Total ({total_segment[2]})"
                )
                slide_count += 1
    return prs


def add_combined_slides(
    prs: Presentation,
    combined_data: List[Tuple[str, List[str], List[Tuple[str, List[float], int]]]],
    group_audience_names: set = None
) -> Presentation:
    """Add combined slides to the presentation."""
    group_audience_names = group_audience_names or set()

    for title, categories, segments in combined_data:
        series_list = []
        footer_parts = ["Source: OnePulse"]

        for label, vals, n_resp in segments:
            series_list.append((label, vals))
            footer_parts.append(f"{label} ({n_resp})")

        # Main combined slide
        slide, chart = create_chart_slide(prs, categories, series_list)
        apply_chart_title(chart, title)
        add_source_footer(slide, prs, ', '.join(footer_parts))

        # Only add individual segment slides for all-segments charts
        non_total_segments = [s for s in segments if s[0] != "Total"]
        is_group_chart = " - " in title
        is_individual_chart = title.endswith(")")

        if len(non_total_segments) > 1 and not is_group_chart and not is_individual_chart:
            total_segment = next((s for s in segments if s[0] == "Total"), None)
            if total_segment:
                for label, vals, n_resp in segments:
                    if label != "Total" and label not in group_audience_names:
                        segment_series = [
                            (total_segment[0], total_segment[1]),
                            (label, vals)
                        ]
                        slide, chart = create_chart_slide(prs, categories, segment_series)
                        apply_chart_title(chart, title)
                        add_source_footer(
                            slide,
                            prs,
                            f"Source: OnePulse, {label} ({n_resp}), Total ({total_segment[2]})"
                        )

    return prs

def add_cover_and_methodology_slides(
    prs: Presentation,
    questions_data: List[Tuple[str, str, List[str]]]
) -> Presentation:
    """
    Add cover and methodology slides to the beginning of the presentation.
    
    Args:
        prs: Presentation object
        questions_data: List of (question_id, question_text, response_options) tuples
    
    Returns:
        Updated presentation
    """
    # Load the cover slides template
    cover_prs = Presentation('cover_slides_template.pptx')
    
    # Get the source slides
    source_slides = list(cover_prs.slides)
    
    # Create new slides at the beginning
    for source_slide in reversed(source_slides):  # Reverse to maintain order
        # Get the layout from the source slide
        layout = source_slide.slide_layout
        
        # Create new slide with the same layout
        new_slide = prs.slides.add_slide(layout)
        
        # Copy shapes from source to new slide
        for shape in source_slide.shapes:
            # Get the shape's position and size
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            
            # Copy the shape based on its type
            if shape.has_text_frame:
                # Create a new text box
                new_shape = new_slide.shapes.add_textbox(left, top, width, height)
                # Copy the text
                new_shape.text_frame.text = shape.text_frame.text
            elif shape.has_chart:
                # Skip charts for now
                continue
            else:
                # For other shapes, just copy the basic properties
                new_shape = new_slide.shapes.add_shape(
                    shape.shape_type,
                    left, top, width, height
                )
    
    return prs

def _question_number_from_text(text: str) -> Optional[int]:
    match = re.match(r'^\s*Q\((\d+)\)', str(text))
    return int(match.group(1)) if match else None


def _normalise_question_text(text: str) -> str:
    return clean_chart_title(str(text)).strip().lower()


def _extract_question_summary(raw_df, raw_audience_data):
    """Build the questionnaire summary directly from the OnePulse export.

    This keeps the original question numbering, includes open-ended questions,
    and identifies appended screener questions from export structure rather
    than broad wording such as "which of the following".
    """
    if raw_df is None:
        return None

    question_ids = set()
    for col in raw_df.columns:
        match = re.match(r'^Q\((\d+)(?:_[^)]+)?\)', str(col))
        if match:
            question_ids.add(int(match.group(1)))

    if not question_ids:
        return None

    # Reuse processed category ordering where possible (e.g. recognised scales).
    processed_lookup = {}
    for title, categories, _ in raw_audience_data or []:
        processed_lookup[_normalise_question_text(title)] = list(categories)

    questions = []
    for q_id in sorted(question_ids):
        prefix = f'Q({q_id})'
        multi_cols = [
            col for col in raw_df.columns
            if re.match(rf'^Q\({q_id}_[^)]+\)', str(col))
        ]
        main_cols = [
            col for col in raw_df.columns
            if str(col).startswith(prefix)
            and not re.match(rf'^Q\({q_id}_[^)]+\)', str(col))
            and 'Comments' not in str(col)
            and 'Sentiment' not in str(col)
        ]

        if multi_cols:
            first_col = str(multi_cols[0])
            match = re.search(r'\[Question:\s*(.*?)\]\s*$', first_col)
            question_text = match.group(1).strip() if match else clean_chart_title(first_col)
            categories = []
            for col in multi_cols:
                label = str(col).split('[', 1)[0].strip()
                label = label.split(')', 1)[1].strip() if ')' in label else label
                categories.append(label)
            is_open_ended = False
            valid_response_count = len(raw_df)
            unique_answer_count = None
        elif main_cols:
            col = main_cols[0]
            question_text = clean_chart_title(str(col))
            series = raw_df[col].dropna()
            valid_response_count = len(series)
            unique_answer_count = series.astype(str).nunique()

            # Long, mostly-unique verbatim answers are treated as open-ended.
            total = len(series)
            unique_ratio = (unique_answer_count / total) if total else 0
            avg_length = series.astype(str).str.len().mean() if total else 0
            is_open_ended = bool(
                total > 0
                and unique_answer_count > 20
                and unique_ratio > 0.5
                and avg_length > 20
            )

            if is_open_ended:
                categories = []
            else:
                processed_categories = processed_lookup.get(_normalise_question_text(question_text))
                if processed_categories:
                    categories = processed_categories
                else:
                    # Preserve first-seen response order as a sensible fallback.
                    categories = list(dict.fromkeys(series.astype(str).tolist()))
        else:
            continue

        questions.append({
            'q_id': q_id,
            'text': question_text,
            'categories': categories,
            'is_open_ended': is_open_ended,
            'valid_response_count': valid_response_count,
            'unique_answer_count': unique_answer_count,
            'is_screener': False,
        })

    # Explicit labels are always respected.
    for q in questions:
        lower_text = q['text'].lower()
        if 'screener' in lower_text or 'screened in' in lower_text:
            q['is_screener'] = True

    # OnePulse appends screeners after the main questionnaire. A common export
    # pattern is that only the qualifying answer remains, so the final question
    # is unanimous across all exported respondents. Mark only trailing unanimous
    # closed questions; this avoids misclassifying normal questions based on wording.
    total_rows = len(raw_df)
    for q in reversed(questions):
        if q['is_screener']:
            continue
        is_unanimous_trailing_gate = (
            not q['is_open_ended']
            and q['unique_answer_count'] == 1
            and q['valid_response_count'] == total_rows
        )
        if is_unanimous_trailing_gate:
            q['is_screener'] = True
        else:
            break

    return questions


def _add_questions_summary_slides(prs, raw_audience_data, raw_df=None):
    """Add one or more questionnaire summary slides."""
    question_summary = _extract_question_summary(raw_df, raw_audience_data)

    if question_summary is None:
        # Backwards-compatible fallback for callers that do not pass raw_df.
        question_summary = []
        for idx, (q_title, categories, _) in enumerate(raw_audience_data or [], start=1):
            lower_text = str(q_title).lower()
            question_summary.append({
                'q_id': _question_number_from_text(q_title) or idx,
                'text': clean_chart_title(q_title),
                'categories': list(categories),
                'is_open_ended': False,
                'is_screener': ('screener' in lower_text or 'screened in' in lower_text),
            })

    audience_labels = []
    if raw_audience_data:
        for label, _, n_resp in raw_audience_data[0][2]:
            audience_labels.append(f"{label} ({n_resp})")

    main_questions = [q for q in question_summary if not q['is_screener']]
    screener_questions = [q for q in question_summary if q['is_screener']]

    # Keep the slide readable: three questionnaire items per summary slide.
    # Screeners are added to the final summary slide when space allows.
    chunks = [main_questions[i:i + 3] for i in range(0, len(main_questions), 3)] or [[]]
    if screener_questions and len(chunks[-1]) >= 3:
        chunks.append([])

    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    for slide_idx, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(blank_layout)
        clear_text(slide)

        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.45), Inches(11.5), Inches(0.5))
        p = title_box.text_frame.paragraphs[0]
        p.text = "Questions and audience" if slide_idx == 0 else "Questions and audience (continued)"
        p.font.name = "Bally Thrill"
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(0, 0, 0)

        body_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.2), Inches(11.5), Inches(5.8))
        tf = body_box.text_frame
        tf.clear()
        tf.word_wrap = True
        first_para = True

        if slide_idx == 0 and audience_labels:
            para = tf.paragraphs[0]
            para.text = "Audience: " + ", ".join(audience_labels)
            para.font.name = "Bally Thrill"
            para.font.size = Pt(12)
            para.font.bold = True
            para.font.color.rgb = RGBColor(0, 0, 0)
            para.space_after = Pt(12)
            first_para = False

        for q in chunk:
            para = tf.paragraphs[0] if first_para else tf.add_paragraph()
            first_para = False
            para.text = f"{q['q_id']}. {q['text']}"
            para.level = 0
            para.font.name = "Bally Thrill"
            para.font.size = Pt(12)
            para.font.bold = True
            para.font.color.rgb = RGBColor(0, 0, 0)
            para.space_before = Pt(8)
            para.space_after = Pt(2)

            para = tf.add_paragraph()
            if q['is_open_ended']:
                para.text = "Response: Open-ended"
            else:
                para.text = "Response options: " + ", ".join(q['categories'])
            para.level = 1
            para.font.name = "Bally Thrill"
            para.font.size = Pt(11)
            para.font.bold = False
            para.font.color.rgb = RGBColor(0, 0, 0)
            para.space_after = Pt(9)

        is_final_summary_slide = slide_idx == len(chunks) - 1
        if is_final_summary_slide and screener_questions:
            para = tf.paragraphs[0] if first_para else tf.add_paragraph()
            first_para = False
            para.text = "SCREENER:"
            para.level = 0
            para.font.name = "Bally Thrill"
            para.font.size = Pt(12)
            para.font.bold = True
            para.font.color.rgb = RGBColor(0, 0, 0)
            para.space_before = Pt(12)
            para.space_after = Pt(2)

            for q in screener_questions:
                para = tf.add_paragraph()
                para.text = q['text']
                para.level = 0
                para.font.name = "Bally Thrill"
                para.font.size = Pt(12)
                para.font.bold = True
                para.font.color.rgb = RGBColor(0, 0, 0)
                para.space_after = Pt(2)

                para = tf.add_paragraph()
                para.text = "Screened in: " + ", ".join(q['categories'])
                para.level = 0
                para.font.name = "Bally Thrill"
                para.font.size = Pt(11)
                para.font.bold = True
                para.font.color.rgb = RGBColor(0, 0, 0)
                para.space_after = Pt(8)

    return prs


def generate_presentation(
    raw_audience_data: List[Tuple[str, List[str], List[Tuple[str, List[float], int]]]],
    combined_data: List[Tuple[str, List[str], List[Tuple[str, List[float], int]]]],
    output_path: Optional[str] = None,
    group_audience_names: set = None,
    export_type: str = "full",
    audience_defs: dict = None,
    raw_df=None
) -> None:
    """
    Generate the complete PowerPoint presentation using raw audience and combined data.
    """
    prs = Presentation(config.TEMPLATE_PATH)

    # Remove any slides already inside the template file
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    # Slide 1: title slide from template
    title_layout = prs.slide_layouts[0]
    s1 = prs.slides.add_slide(title_layout)

    # Replace title slide placeholder text
    text_shapes = [shape for shape in s1.shapes if shape.has_text_frame]
    text_shapes = sorted(text_shapes, key=lambda shape: shape.top)

    if len(text_shapes) >= 1:
        text_shapes[0].text_frame.text = "OnePulse Survey"
        set_text_style(text_shapes[0], font_size=40)

    if len(text_shapes) >= 2:
        text_shapes[1].text_frame.text = "Customer Research Team"
        set_text_style(text_shapes[1], font_size=24)

    # Questionnaire summary slides. Use the raw export when available so the
    # summary reflects the actual questionnaire, including open-ended questions.
    prs = _add_questions_summary_slides(prs, raw_audience_data, raw_df=raw_df)

    # Validate export_type parameter
    if export_type not in ["full", "condensed"]:
        raise ValueError(f"Invalid export_type: {export_type}. Must be 'full' or 'condensed'")

    # Use appropriate function based on export type
    if export_type == "full":
        prs = add_raw_audience_slides(prs, raw_audience_data)
        prs = add_combined_slides_full_export(
            prs,
            combined_data,
            group_audience_names=group_audience_names
        )
    else:
        prs = add_combined_slides_condensed_export(
            prs,
            combined_data,
            group_audience_names=group_audience_names,
            audience_defs=audience_defs,
            raw_audience_data=raw_audience_data
        )

    # Save presentation with appropriate filename
    if output_path:
        output_file = output_path
    else:
        base_name = config.DEFAULT_OUTPUT_PPTX.replace('.pptx', '')
        output_file = f"{base_name}_{export_type}.pptx"

    prs.save(output_file)
    print(f"Presentation saved to: {output_file}")
