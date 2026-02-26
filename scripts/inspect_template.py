import os
from pptx import Presentation

# Get project root
project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

def inspect_slide(slide, slide_index):
    print(f"\nSlide {slide_index + 1}:")
    print("-" * 50)
    
    for shape in slide.shapes:
        print(f"\nShape: {shape.name if hasattr(shape, 'name') else 'Unnamed'}")
        print(f"Type: {shape.shape_type}")
        print(f"ID: {shape.shape_id}")
        
        if shape.has_text_frame:
            print("Text Frame Properties:")
            print(f"Text: {shape.text}")
            print(f"Number of paragraphs: {len(shape.text_frame.paragraphs)}")
            for i, para in enumerate(shape.text_frame.paragraphs):
                print(f"  Paragraph {i + 1}: {para.text}")
                if para.font:
                    print(f"    Font: {para.font.name}")
                    print(f"    Size: {para.font.size}")

def main():
    # Load the template
    template_path = os.path.join(project_root, 'cover_slides_template.pptx')
    prs = Presentation(template_path)
    
    # Inspect each slide
    for i, slide in enumerate(prs.slides):
        inspect_slide(slide, i)

if __name__ == '__main__':
    main() 