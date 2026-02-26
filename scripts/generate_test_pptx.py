#!/usr/bin/env python3
"""
Generate the test PPTX file for inspection.
"""

import os
import sys
import copy
# Add parent directory to path to access src module
project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if project_root not in sys.path:
    sys.path.insert(0, project_root)

from src.data_loader import load_file
from src.main import process_data
from src.ppt_generator import generate_presentation

def generate_test_pptx():
    """Generate the test PPTX file for inspection"""
    print("=== Generating Test PPTX File ===")
    
    # Load test data (same as tests)
    test_csv = os.path.join(project_root, 'tests', 'data', 'test.csv')
    test_data = load_file(test_csv)
    print(f"Loaded test data from: {test_csv}")
    
    # Use EXACTLY the same audience definitions as tests
    audience_defs = {
        "Men": {"Gender": ["Male"]},
        "Women": {"Gender": ["Female"]},
        "Young Adults": {"Age range": [18, 19, 20, 21, 22, 23, 24, 25, 26, 27, 28, 29, 30, 31, 32, 33, 34]},
        "Older Adults": {"Age range": [35, 36, 37, 38, 39, 40, 41, 42, 43, 45, 46, 49, 52, 56, 59, 61, 62, 87]},
        "__groups__": [
            {
                "name": "Gender",
                "audiences": ["Men", "Women"]
            }
        ]
    }
    
    print(f"\n=== Audience Definitions ===")
    print(f"Groups: {[g['name'] for g in audience_defs['__groups__']]}")
    print(f"Audiences: {list(audience_defs.keys())}")
    print(f"Gender group members: {audience_defs['__groups__'][0]['audiences']}")
    
    # Make a copy for process_data (which will mutate it by removing __groups__)
    audience_defs_copy = copy.deepcopy(audience_defs)
    
    # Process the data
    print(f"\n=== Processing data... ===")
    raw_audience_data, combined_data, group_audience_names = process_data(test_data, audience_defs=audience_defs_copy)
    
    print(f"Raw audience data: {len(raw_audience_data)} questions")
    print(f"Combined data: {len(combined_data)} charts")
    print(f"Group audience names: {group_audience_names}")
    
    # Generate presentation
    output_path = os.path.join(project_root, "test_condensed_export.pptx")
    print(f"\n=== Generating presentation: {output_path} ===")
    
    generate_presentation(
        raw_audience_data, 
        combined_data, 
        output_path,
        export_type="condensed",
        audience_defs=audience_defs  # Pass original with __groups__
    )
    
    # Check the output
    if os.path.exists(output_path):
        from pptx import Presentation
        prs = Presentation(output_path)
        print(f"✓ Generated presentation: {output_path}")
        print(f"  - File size: {os.path.getsize(output_path)} bytes")
        print(f"  - Slides: {len(prs.slides)}")
        
        # Show slide titles
        print(f"  - Slide titles:")
        for i, slide in enumerate(prs.slides):
            title = f"Slide {i+1}"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    title = shape.text.strip()
                    break
            print(f"    {i+1}: {title}")
        
        print(f"\n=== File Ready for Inspection ===")
        print(f"The file '{output_path}' is ready for you to open and inspect.")
        print(f"It should contain exactly what the tests expect: 11 slides with group and ungrouped audience charts.")
    else:
        print(f"✗ Failed to generate presentation")

if __name__ == "__main__":
    generate_test_pptx() 