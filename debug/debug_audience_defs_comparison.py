#!/usr/bin/env python3
"""
Debug script to compare different audience definition formats.
This will help us identify what's different between the app and our tests.
"""

import os
import sys
import copy
# Add parent directory to path to access src module
project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if project_root not in sys.path:
    sys.path.insert(0, project_root)

from src.data_loader import load_file
from src.main import process_data
from src.ppt_generator import generate_presentation

def test_audience_defs_comparison():
    """Compare different audience definition formats"""
    print("=== Debug: Audience Definitions Comparison ===")
    
    # Load test data
    test_csv = os.path.join(project_root, 'tests', 'data', 'test.csv')
    test_data = load_file(test_csv)
    print(f"Loaded test data from: {test_csv}")
    
    # Test different audience definition formats
    
    # Format 1: Integer age values (like tests)
    audience_defs_1 = {
        "Men": {"Gender": ["Male"]},
        "Women": {"Gender": ["Female"]},
        "Young Adults": {"Age range": [18, 19, 20, 21, 22, 23, 24, 25, 26, 27, 28, 29, 30, 31, 32, 33, 34]},
        "Older Adults": {"Age range": [35, 36, 37, 38, 39, 40, 41, 42, 43, 45, 46, 49, 52, 56, 59, 61, 62, 87]},
        "__groups__": [
            {
                "name": "Gender",
                "audiences": ["Men", "Women"]
            }
        ]
    }
    
    # Format 2: String age values (like app)
    audience_defs_2 = {
        "Men": {"Gender": ["Male"]},
        "Women": {"Gender": ["Female"]},
        "Young Adults": {"Age range": ["18", "19", "20", "21", "22", "23", "24", "25", "26", "27", "28", "29", "30", "31", "32", "33", "34"]},
        "Older Adults": {"Age range": ["35", "36", "37", "38", "39", "40", "41", "42", "43", "45", "46", "49", "52", "56", "59", "61", "62", "87"]},
        "__groups__": [
            {
                "name": "Gender",
                "audiences": ["Men", "Women"]
            }
        ]
    }
    
    # Format 3: Different audience names (like app might be using)
    audience_defs_3 = {
        "Male": {"Gender": ["Male"]},
        "Female": {"Gender": ["Female"]},
        "Young Adults": {"Age range": ["18", "19", "20", "21", "22", "23", "24", "25", "26", "27", "28", "29", "30", "31", "32", "33", "34"]},
        "Older Adults": {"Age range": ["35", "36", "37", "38", "39", "40", "41", "42", "43", "45", "46", "49", "52", "56", "59", "61", "62", "87"]},
        "__groups__": [
            {
                "name": "Gender",
                "audiences": ["Male", "Female"]
            }
        ]
    }
    
    formats = [
        ("Integer age values (tests)", audience_defs_1),
        ("String age values (app)", audience_defs_2),
        ("Different audience names", audience_defs_3)
    ]
    
    for format_name, audience_defs in formats:
        print(f"\n=== Testing Format: {format_name} ===")
        
        # Make a copy for process_data
        audience_defs_copy = copy.deepcopy(audience_defs)
        
        # Process the data
        raw_audience_data, combined_data, group_audience_names = process_data(test_data, audience_defs=audience_defs_copy)
        
        print(f"Raw audience data: {len(raw_audience_data)} questions")
        print(f"Combined data: {len(combined_data)} charts")
        print(f"Group audience names: {group_audience_names}")
        
        # Count chart types
        group_charts = []
        individual_charts = []
        all_segments_charts = []
        
        for title, categories, segments in combined_data:
            is_group_chart = " - " in title
            is_individual_chart = title.endswith(")")
            
            if is_group_chart:
                group_charts.append((title, [s[0] for s in segments]))
            elif is_individual_chart:
                individual_charts.append((title, [s[0] for s in segments]))
            else:
                all_segments_charts.append((title, [s[0] for s in segments]))
        
        print(f"Group charts: {len(group_charts)}")
        print(f"Individual charts: {len(individual_charts)}")
        print(f"All segments charts: {len(all_segments_charts)}")
        
        # Test PowerPoint generation
        output_path = f"debug_comparison_{format_name.replace(' ', '_').replace('(', '').replace(')', '')}.pptx"
        
        try:
            generate_presentation(
                raw_audience_data, 
                combined_data, 
                output_path,
                export_type="condensed",
                audience_defs=audience_defs
            )
            
            if os.path.exists(output_path):
                from pptx import Presentation
                prs = Presentation(output_path)
                print(f"✓ Generated {len(prs.slides)} slides")
                
                # Check for group slides
                group_slides = 0
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and "Men & Women" in shape.text:
                            group_slides += 1
                            break
                
                print(f"  - Group slides found: {group_slides}")
                
                # Clean up
                os.unlink(output_path)
            else:
                print("✗ Failed to generate presentation")
                
        except Exception as e:
            print(f"✗ Error: {e}")
    
    print(f"\n=== Analysis ===")
    print("If all formats produce the same results, the issue is elsewhere.")
    print("If one format produces different results, that's the issue.")

if __name__ == "__main__":
    test_audience_defs_comparison() 