#!/usr/bin/env python3
"""
Debug script to check what data the app is actually using.
"""

import os
import sys
import json
import copy
# Add parent directory to path to access src module
project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if project_root not in sys.path:
    sys.path.insert(0, project_root)

from src.data_loader import load_file
from src.main import process_data
from src.ppt_generator import generate_presentation

def debug_app_data_check():
    """Check what data the app is actually using"""
    print("=== Debug: App Data Check ===")
    
    # Check if the app's audience_segments.json exists
    app_json_path = os.path.join(project_root, "src", "audience_segments.json")
    if os.path.exists(app_json_path):
        print(f"✓ Found app's audience_segments.json: {app_json_path}")
        with open(app_json_path, 'r') as f:
            app_json = json.load(f)
        print(f"App JSON has __groups__: {'__groups__' in app_json}")
        if '__groups__' in app_json:
            print(f"App groups: {app_json['__groups__']}")
    else:
        print(f"✗ App's audience_segments.json not found: {app_json_path}")
        return
    
    # Check what data file the app might be using
    possible_data_files = [
        os.path.join(project_root, 'tests', 'data', 'test.csv'),
        os.path.join(project_root, 'tests', 'data', 'test.xlsx'),
        os.path.join(project_root, 'survey_data', 'test_data_for_app.xlsx'),
        os.path.join(project_root, 'survey_data', 'Problem_file.csv'),
    ]
    
    print(f"\n=== Checking Data Files ===")
    for data_file in possible_data_files:
        if os.path.exists(data_file):
            print(f"✓ Found data file: {data_file}")
            try:
                test_data = load_file(data_file)
                print(f"  - Shape: {test_data.shape}")
                print(f"  - Columns: {len(test_data.columns)}")
                print(f"  - Sample columns: {list(test_data.columns)[:5]}")
                
                # Test with this data
                print(f"  - Testing with this data...")
                audience_defs_copy = copy.deepcopy(app_json)
                
                raw_audience_data, combined_data, group_audience_names = process_data(test_data, audience_defs=audience_defs_copy)
                
                print(f"    Raw audience data: {len(raw_audience_data)} questions")
                print(f"    Combined data: {len(combined_data)} charts")
                print(f"    Group audience names: {group_audience_names}")
                
                # Count chart types
                group_charts = []
                individual_charts = []
                all_segments_charts = []
                
                for title, categories, segments in combined_data:
                    is_group_chart = " - " in title
                    is_individual_chart = title.endswith(")")
                    
                    if is_group_chart:
                        group_charts.append((title, [s[0] for s in segments]))
                    elif is_individual_chart:
                        individual_charts.append((title, [s[0] for s in segments]))
                    else:
                        all_segments_charts.append((title, [s[0] for s in segments]))
                
                print(f"    Group charts: {len(group_charts)}")
                print(f"    Individual charts: {len(individual_charts)}")
                print(f"    All segments charts: {len(all_segments_charts)}")
                
                # Generate presentation
                output_path = f"debug_app_data_check_{os.path.basename(data_file).replace('.', '_')}.pptx"
                
                generate_presentation(
                    raw_audience_data, 
                    combined_data, 
                    output_path,
                    export_type="condensed",
                    audience_defs=app_json
                )
                
                if os.path.exists(output_path):
                    from pptx import Presentation
                    prs = Presentation(output_path)
                    print(f"    ✓ Generated {len(prs.slides)} slides")
                    
                    # Count group slides
                    group_slides = 0
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and "Men & Women" in shape.text:
                                group_slides += 1
                                break
                    
                    print(f"    Group slides found: {group_slides}")
                    
                    # Clean up
                    os.unlink(output_path)
                    
                    # Check if this matches the user's issue
                    if len(prs.slides) == 8:
                        print(f"    ❌ MATCHES USER'S ISSUE: 8 slides (missing group slides)")
                    elif len(prs.slides) == 11:
                        print(f"    ✅ EXPECTED RESULT: 11 slides (with group slides)")
                    else:
                        print(f"    ⚠️  UNEXPECTED: {len(prs.slides)} slides")
                        
                else:
                    print(f"    ✗ Failed to generate presentation")
                    
            except Exception as e:
                print(f"  - Error processing {data_file}: {e}")
        else:
            print(f"✗ Data file not found: {data_file}")
    
    print(f"\n=== Analysis ===")
    print("If any data file produces 8 slides, that's the data the app is using.")
    print("If all data files produce 11 slides, the issue is elsewhere in the app.")

if __name__ == "__main__":
    debug_app_data_check() 