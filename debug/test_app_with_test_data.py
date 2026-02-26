#!/usr/bin/env python3
"""
Test script to run the actual app with the test data file.
This will help us understand what the real output looks like vs. the test expectations.
"""

import os
import sys
import copy
# Add parent directory to path to access src module
project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if project_root not in sys.path:
    sys.path.insert(0, project_root)

from src.data_loader import load_file
from src.main import process_data
from src.ppt_generator import generate_presentation
from src import config

def test_app_with_test_data():
    """Test the app with the same test data and audience definitions as the tests"""
    print("=== Testing App with Test Data ===")
    
    # Load test data (same as tests)
    test_csv = os.path.join(project_root, 'tests', 'data', 'test.csv')
    test_data = load_file(test_csv)
    print(f"Loaded test data from: {test_csv}")
    print(f"Test data shape: {test_data.shape}")
    print(f"Test data columns: {list(test_data.columns)}")
    
    # Use the same audience definitions as TestExportTypes
    audience_defs = {
        "Men": {"Gender": ["Male"]},
        "Women": {"Gender": ["Female"]},
        "Young Adults": {"Age range": [18, 19, 20, 21, 22, 23, 24, 25, 26, 27, 28, 29, 30, 31, 32, 33, 34]},
        "Older Adults": {"Age range": [35, 36, 37, 38, 39, 40, 41, 42, 43, 45, 46, 49, 52, 56, 59, 61, 62, 87]},
        "__groups__": [
            {
                "name": "Gender",
                "audiences": ["Men", "Women"]
            },
            {
                "name": "Age",
                "audiences": ["Young Adults", "Older Adults"]
            }
        ]
    }
    
    print(f"\nAudience definitions:")
    print(f"- Groups: {[g['name'] for g in audience_defs['__groups__']]}")
    print(f"- Audiences: {list(audience_defs.keys())}")
    
    # Make a copy for process_data (which will mutate it by removing __groups__)
    audience_defs_copy = copy.deepcopy(audience_defs)
    
    # Process the data
    print(f"\nProcessing data...")
    raw_audience_data, combined_data, group_audience_names = process_data(test_data, audience_defs=audience_defs_copy)
    
    print(f"Raw audience data: {len(raw_audience_data)} questions")
    print(f"Combined data: {len(combined_data)} charts")
    print(f"Group audience names: {group_audience_names}")
    
    # Test both export types
    for export_type in ["full", "condensed"]:
        print(f"\n=== Testing {export_type.upper()} Export ===")
        
        # Generate filename with export type
        output_filename = f"test_data_{export_type}.pptx"
        output_path = os.path.join(project_root, "exports", output_filename)
        
        # Ensure exports directory exists
        os.makedirs("exports", exist_ok=True)
        
        # Generate presentation
        print(f"Generating {export_type} presentation...")
        generate_presentation(
            raw_audience_data, 
            combined_data, 
            output_path,
            export_type=export_type,
            audience_defs=audience_defs  # Pass original with __groups__
        )
        
        # Check the output
        if os.path.exists(output_path):
            from pptx import Presentation
            prs = Presentation(output_path)
            print(f"✓ Generated {export_type} presentation: {output_path}")
            print(f"  - Slides: {len(prs.slides)}")
            
            # Show slide titles
            print(f"  - Slide titles:")
            for i, slide in enumerate(prs.slides):
                title = f"Slide {i+1}"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        title = shape.text.strip()
                        break
                print(f"    {i+1}: {title}")
        else:
            print(f"✗ Failed to generate {export_type} presentation")
    
    print(f"\n=== Test Complete ===")
    print(f"Check the 'exports' directory for the generated files:")
    print(f"- test_data_full.pptx")
    print(f"- test_data_condensed.pptx")

if __name__ == "__main__":
    test_app_with_test_data() 