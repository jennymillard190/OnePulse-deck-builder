#!/usr/bin/env python3
"""
Debug script to test condensed export issues:
1. No slides when no audiences are added
2. Audiences not grouped together properly
"""

import os
import sys
# Add parent directory to path to access src module
project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if project_root not in sys.path:
    sys.path.insert(0, project_root)

from src.data_loader import load_file
from src.main import process_data
from src.ppt_generator import generate_presentation

def test_condensed_no_audiences():
    """Test condensed export with no audiences"""
    print("=== Testing Condensed Export with NO Audiences ===")
    
    # Load test data
    test_csv = os.path.join(project_root, 'tests', 'data', 'test.csv')
    test_data = load_file(test_csv)
    
    # Process data with no audiences
    raw_audience_data, combined_data, group_audience_names = process_data(test_data)
    
    print(f"Raw audience data items: {len(raw_audience_data)}")
    print(f"Combined data items: {len(combined_data)}")
    
    # Generate condensed presentation
    output_path = "debug_condensed_no_audiences.pptx"
    try:
        generate_presentation(
            raw_audience_data, 
            combined_data, 
            output_path,
            export_type="condensed"
        )
        
        # Check if file was created and has content
        if os.path.exists(output_path):
            from pptx import Presentation
            prs = Presentation(output_path)
            print(f"Generated presentation has {len(prs.slides)} slides")
            
            # List slide titles
            for i, slide in enumerate(prs.slides):
                print(f"  Slide {i}: {[shape.text for shape in slide.shapes if hasattr(shape, 'text')]}")
        else:
            print("No presentation file was created!")
            
    except Exception as e:
        print(f"Error generating presentation: {e}")
    finally:
        if os.path.exists(output_path):
            os.unlink(output_path)

def test_condensed_with_groups():
    """Test condensed export with grouped audiences"""
    print("\n=== Testing Condensed Export with Grouped Audiences ===")
    
    # Load test data
    test_csv = os.path.join(project_root, 'tests', 'data', 'test.csv')
    test_data = load_file(test_csv)
    
    # Define audiences with groups
    audience_defs = {
        "Men": {"Gender": ["Male"]},
        "Women": {"Gender": ["Female"]},
        "Young Adults": {"Age range": [18, 19, 20, 21, 22, 23, 24, 25, 26, 27, 28, 29, 30, 31, 32, 33, 34]},
        "Older Adults": {"Age range": [35, 36, 37, 38, 39, 40, 41, 42, 43, 45, 46, 49, 52, 56, 59, 61, 62, 87]},
        "__groups__": [
            {
                "name": "Gender",
                "audiences": ["Men", "Women"]
            },
            {
                "name": "Age",
                "audiences": ["Young Adults", "Older Adults"]
            }
        ]
    }
    
    # Process data with audiences
    raw_audience_data, combined_data, group_audience_names = process_data(test_data, audience_defs=audience_defs)
    
    print(f"Raw audience data items: {len(raw_audience_data)}")
    print(f"Combined data items: {len(combined_data)}")
    print(f"Group audience names: {group_audience_names}")
    
    # List combined data structure
    print("\nCombined data structure:")
    for i, (title, categories, segments) in enumerate(combined_data):
        print(f"  {i}: Title='{title}'")
        print(f"     Categories: {len(categories)}")
        print(f"     Segments: {[s[0] for s in segments]}")
    
    # Generate condensed presentation
    output_path = "debug_condensed_with_groups.pptx"
    try:
        generate_presentation(
            raw_audience_data, 
            combined_data, 
            output_path,
            export_type="condensed",
            audience_defs=audience_defs
        )
        
        # Check if file was created and has content
        if os.path.exists(output_path):
            from pptx import Presentation
            prs = Presentation(output_path)
            print(f"\nGenerated presentation has {len(prs.slides)} slides")
            
            # List slide titles
            for i, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        slide_texts.append(shape.text)
                print(f"  Slide {i}: {slide_texts}")
        else:
            print("No presentation file was created!")
            
    except Exception as e:
        print(f"Error generating presentation: {e}")
    finally:
        if os.path.exists(output_path):
            os.unlink(output_path)

if __name__ == "__main__":
    test_condensed_no_audiences()
    test_condensed_with_groups() 