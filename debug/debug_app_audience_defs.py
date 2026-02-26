#!/usr/bin/env python3
"""
Debug script to test the app's audience definitions processing.
This will help us understand why group slides aren't being generated in the app.
"""

import os
import sys
import copy
import json
# Add parent directory to path to access src module
project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if project_root not in sys.path:
    sys.path.insert(0, project_root)

from src.data_loader import load_file
from src.main import process_data
from src.ppt_generator import generate_presentation

def test_app_audience_defs():
    """Test the app's audience definitions processing"""
    print("=== Testing App Audience Definitions Processing ===")
    
    # Load test data (same as tests)
    test_csv = os.path.join(project_root, 'tests', 'data', 'test.csv')
    test_data = load_file(test_csv)
    print(f"Loaded test data from: {test_csv}")
    print(f"Test data shape: {test_data.shape}")
    
    # Simulate the app's audience definitions (Gender group + ungrouped Young/Old)
    audience_defs = {
        "Male": {"Gender": ["Male"]},
        "Female": {"Gender": ["Female"]},
        "Young": {"Age range": [18, 19, 20, 21, 22, 23, 24, 25, 26, 27, 28, 29, 30, 31, 32, 33, 34]},
        "Old": {"Age range": [35, 36, 37, 38, 39, 40, 41, 42, 43, 45, 46, 49, 52, 56, 59, 61, 62, 87]},
        "__groups__": [
            {
                "name": "Gender",
                "audiences": ["Male", "Female"]
            }
        ]
    }
    
    print(f"\n=== Original Audience Definitions ===")
    print(f"Groups: {[g['name'] for g in audience_defs['__groups__']]}")
    print(f"Audiences: {list(audience_defs.keys())}")
    print(f"Gender group members: {audience_defs['__groups__'][0]['audiences']}")
    
    # Make a copy for process_data (which will mutate it by removing __groups__)
    audience_defs_copy = copy.deepcopy(audience_defs)
    
    print(f"\n=== Before process_data ===")
    print(f"Copy has __groups__: {'__groups__' in audience_defs_copy}")
    print(f"Copy groups: {audience_defs_copy.get('__groups__', [])}")
    
    # Process the data
    print(f"\n=== Processing data... ===")
    raw_audience_data, combined_data, group_audience_names = process_data(test_data, audience_defs=audience_defs_copy)
    
    print(f"Raw audience data: {len(raw_audience_data)} questions")
    print(f"Combined data: {len(combined_data)} charts")
    print(f"Group audience names: {group_audience_names}")
    
    # Check what's in combined_data
    print(f"\n=== Combined Data Analysis ===")
    group_charts = []
    individual_charts = []
    all_segments_charts = []
    
    for title, categories, segments in combined_data:
        is_group_chart = " - " in title
        is_individual_chart = title.endswith(")")
        
        if is_group_chart:
            group_charts.append((title, [s[0] for s in segments]))
        elif is_individual_chart:
            individual_charts.append((title, [s[0] for s in segments]))
        else:
            all_segments_charts.append((title, [s[0] for s in segments]))
    
    print(f"Group charts ({len(group_charts)}):")
    for title, segments in group_charts:
        print(f"  - {title}: {segments}")
    
    print(f"Individual charts ({len(individual_charts)}):")
    for title, segments in individual_charts:
        print(f"  - {title}: {segments}")
    
    print(f"All segments charts ({len(all_segments_charts)}):")
    for title, segments in all_segments_charts:
        print(f"  - {title}: {segments}")
    
    # Test condensed export
    print(f"\n=== Testing Condensed Export ===")
    output_path = "debug_app_condensed.pptx"
    
    # Generate presentation with original audience_defs (with __groups__)
    generate_presentation(
        raw_audience_data, 
        combined_data, 
        output_path,
        export_type="condensed",
        audience_defs=audience_defs  # Pass original with __groups__
    )
    
    # Check the output
    if os.path.exists(output_path):
        from pptx import Presentation
        prs = Presentation(output_path)
        print(f"✓ Generated condensed presentation: {output_path}")
        print(f"  - Slides: {len(prs.slides)}")
        
        # Show slide titles
        print(f"  - Slide titles:")
        for i, slide in enumerate(prs.slides):
            title = f"Slide {i+1}"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    title = shape.text.strip()
                    break
            print(f"    {i+1}: {title}")
        
        # Clean up
        os.unlink(output_path)
    else:
        print(f"✗ Failed to generate condensed presentation")
    
    print(f"\n=== Analysis ===")
    if len(group_charts) == 0:
        print("❌ PROBLEM: No group charts found in combined_data!")
        print("   This means the data processing is not creating group charts.")
        print("   Check that audience_defs has __groups__ and group members exist in data.")
    else:
        print("✅ Group charts found in combined_data")
    
    expected_group_charts = len(audience_defs['__groups__']) * len(raw_audience_data)
    if len(group_charts) != expected_group_charts:
        print(f"⚠️  WARNING: Expected {expected_group_charts} group charts, got {len(group_charts)}")
    
    print(f"\n=== Test Complete ===")

if __name__ == "__main__":
    test_app_audience_defs() 