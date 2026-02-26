#!/usr/bin/env python3
"""
Debug script to run the test data through the app and see what the actual output looks like.
This will help us understand if the code matches the test expectations.
"""

import os
import sys
# Add parent directory to path to access src module
project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if project_root not in sys.path:
    sys.path.insert(0, project_root)

import copy
from src.data_loader import load_file
from src.main import process_data
from src.ppt_generator import generate_presentation

def test_with_same_audience_defs_as_tests():
    """Test with the same audience definitions as the TestExportTypes class"""
    print("=== Testing with Same Audience Definitions as Tests ===")
    
    # Load test data
    test_csv = os.path.join(project_root, 'tests', 'data', 'test.csv')
    test_data = load_file(test_csv)
    
    # Use the same audience definitions as TestExportTypes
    audience_defs = {
        "Men": {"Gender": ["Male"]},
        "Women": {"Gender": ["Female"]},
        "Young Adults": {"Age range": [18, 19, 20, 21, 22, 23, 24, 25, 26, 27, 28, 29, 30, 31, 32, 33, 34]},
        "Older Adults": {"Age range": [35, 36, 37, 38, 39, 40, 41, 42, 43, 45, 46, 49, 52, 56, 59, 61, 62, 87]},
        "__groups__": [
            {
                "name": "Gender",
                "audiences": ["Men", "Women"]
            },
            {
                "name": "Age",
                "audiences": ["Young Adults", "Older Adults"]
            }
        ]
    }
    
    print(f"Audience definitions: {audience_defs}")
    
    # Pass a deep copy to process_data so original is not mutated
    raw_audience_data, combined_data, group_audience_names = process_data(test_data, audience_defs=copy.deepcopy(audience_defs))
    
    print(f"\nRaw audience data items: {len(raw_audience_data)}")
    print(f"Combined data items: {len(combined_data)}")
    print(f"Group audience names: {group_audience_names}")
    
    # List combined data structure
    print("\nCombined data structure:")
    for i, (title, categories, segments) in enumerate(combined_data):
        print(f"  {i}: Title='{title}'")
        print(f"     Categories: {len(categories)}")
        print(f"     Segments: {[s[0] for s in segments]}")
    
    # Generate condensed presentation
    output_path = "debug_test_data_condensed.pptx"
    try:
        generate_presentation(
            raw_audience_data, 
            combined_data, 
            output_path,
            export_type="condensed",
            audience_defs=audience_defs
        )
        
        # Load the presentation and analyze it
        from pptx import Presentation
        prs = Presentation(output_path)
        
        print(f"\nGenerated presentation has {len(prs.slides)} slides")
        
        # Analyze each slide
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            slide_titles = []
            
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    slide_texts.append(shape.text)
                    # Look for chart titles (they usually contain the question text)
                    if "Question:" in shape.text or "Q(" in shape.text:
                        slide_titles.append(shape.text)
            
            print(f"\n  Slide {i}:")
            print(f"    All text: {slide_texts}")
            if slide_titles:
                print(f"    Chart titles: {slide_titles}")
            
            # Look for group indicators in slide content
            group_indicators = []
            for text in slide_texts:
                if "Men & Women" in text:
                    group_indicators.append("Gender group")
                if "Young Adults & Older Adults" in text:
                    group_indicators.append("Age group")
                if "Source: OnePulse, Men" in text and "Source: OnePulse, Women" in text:
                    group_indicators.append("Individual Men/Women slides")
            
            if group_indicators:
                print(f"    Group indicators found: {group_indicators}")
        
        # Expected vs actual analysis
        print(f"\n=== ANALYSIS ===")
        print(f"Expected slides: 2 cover + (3 questions × (2 groups + 0 ungrouped)) = 8 slides")
        print(f"Actual slides: {len(prs.slides)}")
        
        if len(prs.slides) == 8:
            print("✅ Slide count matches expectation")
        else:
            print(f"❌ Slide count mismatch: expected 8, got {len(prs.slides)}")
        
        # Check for group slides
        gender_group_slides = 0
        age_group_slides = 0
        individual_slides = 0
        
        for i in range(2, len(prs.slides)):  # Skip cover slides
            slide = prs.slides[i]
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    if "Men & Women" in shape.text:
                        gender_group_slides += 1
                    elif "Young Adults & Older Adults" in shape.text:
                        age_group_slides += 1
                    elif "Source: OnePulse, Men" in shape.text or "Source: OnePulse, Women" in shape.text:
                        individual_slides += 1
        
        print(f"\nGroup slides found:")
        print(f"  Gender group slides: {gender_group_slides}")
        print(f"  Age group slides: {age_group_slides}")
        print(f"  Individual audience slides: {individual_slides}")
        
        if gender_group_slides == 3 and age_group_slides == 3 and individual_slides == 0:
            print("✅ Group slides match expectation (3 Gender + 3 Age, no individual)")
        else:
            print("❌ Group slides don't match expectation")
            
    finally:
        if os.path.exists(output_path):
            os.unlink(output_path)

if __name__ == "__main__":
    test_with_same_audience_defs_as_tests() 