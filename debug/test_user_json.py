#!/usr/bin/env python3
"""
Test script for the user's exact JSON to identify the issue.
"""

import os
import sys
import json
import copy
# Add parent directory to path to access src module
project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if project_root not in sys.path:
    sys.path.insert(0, project_root)

from src.data_loader import load_file
from src.main import process_data
from src.ppt_generator import generate_presentation

def test_user_json():
    """Test the user's exact JSON"""
    print("=== Testing User's Exact JSON ===")
    
    # Load test data
    test_csv = os.path.join(project_root, 'tests', 'data', 'test.csv')
    test_data = load_file(test_csv)
    print(f"Loaded test data from: {test_csv}")
    
    # The user's exact JSON (from their description)
    user_json = {
        "Men": {
            "Gender": [
                "Male"
            ]
        },
        "Women": {
            "Gender": [
                "Female"
            ]
        },
        "Young Adults": {
            "Age range": [
                "18",
                "19",
                "20",
                "21",
                "22",
                "23",
                "24",
                "25",
                "26",
                "27",
                "28",
                "29",
                "30",
                "31",
                "32",
                "33",
                "34"
            ]
        },
        "Older Adults": {
            "Age range": [
                "35",
                "36",
                "37",
                "38",
                "39",
                "40",
                "41",
                "42",
                "43",
                "45",
                "46",
                "49",
                "52",
                "56",
                "59",
                "61",
                "62",
                "87"
            ]
        },
        "__groups__": [
            {
                "name": "Gender",
                "audiences": [
                    "Men",
                    "Women"
                ]
            }
        ]
    }
    
    print(f"\n=== User's JSON ===")
    print(json.dumps(user_json, indent=2))
    
    # Make a copy for process_data
    audience_defs_copy = copy.deepcopy(user_json)
    
    # Process the data
    print(f"\n=== Processing data... ===")
    raw_audience_data, combined_data, group_audience_names = process_data(test_data, audience_defs=audience_defs_copy)
    
    print(f"Raw audience data: {len(raw_audience_data)} questions")
    print(f"Combined data: {len(combined_data)} charts")
    print(f"Group audience names: {group_audience_names}")
    
    # Analyze combined data
    print(f"\n=== Combined Data Analysis ===")
    group_charts = []
    individual_charts = []
    all_segments_charts = []
    
    for title, categories, segments in combined_data:
        is_group_chart = " - " in title
        is_individual_chart = title.endswith(")")
        
        if is_group_chart:
            group_charts.append((title, [s[0] for s in segments]))
        elif is_individual_chart:
            individual_charts.append((title, [s[0] for s in segments]))
        else:
            all_segments_charts.append((title, [s[0] for s in segments]))
    
    print(f"Group charts ({len(group_charts)}):")
    for title, segments in group_charts:
        print(f"  - {title}: {segments}")
    
    print(f"Individual charts ({len(individual_charts)}):")
    for title, segments in individual_charts:
        print(f"  - {title}: {segments}")
    
    print(f"All segments charts ({len(all_segments_charts)}):")
    for title, segments in all_segments_charts:
        print(f"  - {title}: {segments}")
    
    # Generate presentation
    print(f"\n=== Generating presentation... ===")
    output_path = "test_user_json.pptx"
    
    generate_presentation(
        raw_audience_data, 
        combined_data, 
        output_path,
        export_type="condensed",
        audience_defs=user_json
    )
    
    # Check the output
    if os.path.exists(output_path):
        from pptx import Presentation
        prs = Presentation(output_path)
        print(f"✓ Generated presentation: {output_path}")
        print(f"  - Slides: {len(prs.slides)}")
        
        # Show slide titles
        print(f"  - Slide titles:")
        for i, slide in enumerate(prs.slides):
            title = f"Slide {i+1}"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    title = shape.text.strip()
                    break
            print(f"    {i+1}: {title}")
        
        # Count group slides
        group_slides = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and "Men & Women" in shape.text:
                    group_slides += 1
                    break
        
        print(f"  - Group slides found: {group_slides}")
        
        # Clean up
        os.unlink(output_path)
    else:
        print(f"✗ Failed to generate presentation")
    
    print(f"\n=== Expected vs Actual ===")
    print(f"Expected slides: 2 cover + 3 questions × (1 group + 2 ungrouped) = 11")
    print(f"Expected group charts: {len(user_json['__groups__']) * len(raw_audience_data)}")
    print(f"Expected individual charts: {len(raw_audience_data) * 2} (3 questions × 2 ungrouped)")

if __name__ == "__main__":
    test_user_json() 