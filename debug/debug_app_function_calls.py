#!/usr/bin/env python3
"""
Debug script that exactly replicates the app's function calls.
"""

import os
import sys
import json
import copy
# Add parent directory to path to access src module
project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if project_root not in sys.path:
    sys.path.insert(0, project_root)

from src.data_loader import load_file
from src.main import process_data
from src.ppt_generator import generate_presentation
from src import config

def debug_app_function_calls():
    """Exactly replicate the app's function calls"""
    print("=== Debug: App Function Calls ===")
    
    # Load test data
    test_csv = os.path.join(project_root, 'tests', 'data', 'test.csv')
    test_data = load_file(test_csv)
    print(f"Loaded test data from: {test_csv}")
    
    # Load the app's JSON
    app_json_path = os.path.join(project_root, "src", "audience_segments.json")
    with open(app_json_path, 'r') as f:
        app_json = json.load(f)
    
    print(f"Loaded app JSON with groups: {app_json.get('__groups__', [])}")
    
    # EXACTLY replicate the app's function calls
    print(f"\n=== Replicating App's Function Calls ===")
    
    # Step 1: Process the data (EXACTLY like the app does)
    print(f"Step 1: Calling process_data...")
    raw_audience_data, combined_data, group_audience_names = process_data(test_data, audience_defs=app_json)
    
    print(f"  Raw audience data: {len(raw_audience_data)} questions")
    print(f"  Combined data: {len(combined_data)} charts")
    print(f"  Group audience names: {group_audience_names}")
    
    # Step 2: Generate presentation (EXACTLY like the app does)
    print(f"\nStep 2: Calling generate_presentation...")
    output_path = os.path.join(project_root, "exports", "test_condensed.pptx")
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    
    # Use the EXACT same parameters as the app
    generate_presentation(
        raw_audience_data, 
        combined_data, 
        output_path,
        export_type="condensed",
        audience_defs=app_json
    )
    
    # Check the output
    if os.path.exists(output_path):
        from pptx import Presentation
        prs = Presentation(output_path)
        print(f"✓ Generated presentation: {output_path}")
        print(f"  - Slides: {len(prs.slides)}")
        
        # Show slide titles
        print(f"  - Slide titles:")
        for i, slide in enumerate(prs.slides):
            title = f"Slide {i+1}"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    title = shape.text.strip()
                    break
            print(f"    {i+1}: {title}")
        
        # Count group slides
        group_slides = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and "Men & Women" in shape.text:
                    group_slides += 1
                    break
        
        print(f"  - Group slides found: {group_slides}")
        
        # Check if this matches what the user sees
        if len(prs.slides) == 8:
            print(f"\n❌ MATCHES USER'S ISSUE: Got 8 slides (missing group slides)")
        elif len(prs.slides) == 11:
            print(f"\n✅ EXPECTED RESULT: Got 11 slides (with group slides)")
        else:
            print(f"\n⚠️  UNEXPECTED: Got {len(prs.slides)} slides")
    else:
        print(f"✗ Failed to generate presentation")
    
    # Now let's check what the app's actual output file contains
    print(f"\n=== Checking App's Actual Output File ===")
    app_output_path = os.path.join(project_root, "exports", "test_condensed.pptx")
    if os.path.exists(app_output_path):
        from pptx import Presentation
        prs = Presentation(app_output_path)
        print(f"App's actual output file: {app_output_path}")
        print(f"  - Slides: {len(prs.slides)}")
        
        # Show slide titles
        print(f"  - Slide titles:")
        for i, slide in enumerate(prs.slides):
            title = f"Slide {i+1}"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    title = shape.text.strip()
                    break
            print(f"    {i+1}: {title}")
        
        # Count group slides
        group_slides = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and "Men & Women" in shape.text:
                    group_slides += 1
                    break
        
        print(f"  - Group slides found: {group_slides}")
    else:
        print(f"App's output file not found: {app_output_path}")

if __name__ == "__main__":
    debug_app_function_calls() 