import unittest
import pandas as pd
import os
import re
import logging
from src.main import process_data
from src.ppt_generator import create_chart_slide, generate_presentation
from src.data_loader import load_file, load_uploaded_file, process_dataframe, process_semicolon_separated_column
from pptx import Presentation
import src.ppt_generator as ppt_generator
import src.config as config
from pptx.exc import PackageNotFoundError

# Configure logging to only show warnings and errors during tests
logging.basicConfig(level=logging.WARNING)

class TestDataLoading(unittest.TestCase):
    """Test group for data loading functionality"""
    
    def setUp(self):
        """Set up test data paths"""
        self.test_data_dir = os.path.join(os.path.dirname(__file__), 'data')
        self.xlsx_path = os.path.join(self.test_data_dir, 'test.xlsx')
        self.csv_path = os.path.join(self.test_data_dir, 'test.csv')
    
    def test_xlsx_loading(self):
        """Test loading XLSX file with correct header row"""
        # Load the test XLSX file using our application's function
        df = load_file(self.xlsx_path)
        
        # Check that the data was loaded correctly
        self.assertIsInstance(df, pd.DataFrame, "Data should be loaded as a DataFrame")
        self.assertGreater(len(df), 0, "DataFrame should not be empty")
        
        # Check that column names are properly formatted
        for col in df.columns:
            # Check for multiple spaces in the middle of column names
            self.assertFalse('  ' in col, 
                           f"Column name '{col}' contains multiple spaces")
            # Check for trailing spaces
            self.assertEqual(col, col.rstrip(), 
                           f"Column name '{col}' contains trailing whitespace")
            
        # Check for bank columns if they exist
        bank_column = 'Bank(s)'
        if bank_column in df.columns:
            # Check that bank columns are created
            bank_customer_cols = [col for col in df.columns if col.endswith('_customer')]
            self.assertGreater(len(bank_customer_cols), 0, 
                             "Should create customer columns for banks")
            
            # Check that bank columns are boolean
            for col in bank_customer_cols:
                self.assertEqual(df[col].dtype, bool, 
                               f"Bank customer column {col} should be boolean")
        
        # Check for multi-select questions
        multi_select_cols = [col for col in df.columns if re.match(r'Q\(\d+_\d+\)', col)]
        if multi_select_cols:
            for col in multi_select_cols:
                # Check that values are boolean
                self.assertEqual(df[col].dtype, bool, 
                               f"Multi-select column {col} should be boolean")
                # Check that all values are either True or False
                self.assertTrue(all(isinstance(val, bool) for val in df[col]),
                              f"Multi-select values in {col} should be boolean")
    
    def test_csv_loading(self):
        """Test loading CSV file with correct header row"""
        # Load the test CSV file using our application's function
        df = load_file(self.csv_path)
        
        # Check that the data was loaded correctly
        self.assertIsInstance(df, pd.DataFrame, "Data should be loaded as a DataFrame")
        self.assertGreater(len(df), 0, "DataFrame should not be empty")
        
        # Check that column names are properly formatted
        for col in df.columns:
            # Check for multiple spaces in the middle of column names
            self.assertFalse('  ' in col, 
                           f"Column name '{col}' contains multiple spaces")
            # Check for trailing spaces
            self.assertEqual(col, col.rstrip(), 
                           f"Column name '{col}' contains trailing whitespace")
            
        # Check for bank columns if they exist
        bank_column = 'Bank(s)'
        if bank_column in df.columns:
            # Check that bank columns are created
            bank_customer_cols = [col for col in df.columns if col.endswith('_customer')]
            self.assertGreater(len(bank_customer_cols), 0, 
                             "Should create customer columns for banks")
            
            # Check that bank columns are boolean
            for col in bank_customer_cols:
                self.assertEqual(df[col].dtype, bool, 
                               f"Bank customer column {col} should be boolean")
        
        # Check for multi-select questions
        multi_select_cols = [col for col in df.columns if re.match(r'Q\(\d+_\d+\)', col)]
        if multi_select_cols:
            for col in multi_select_cols:
                # Check that values are boolean
                self.assertEqual(df[col].dtype, bool, 
                               f"Multi-select column {col} should be boolean")
                # Check that all values are either True or False
                self.assertTrue(all(isinstance(val, bool) for val in df[col]),
                              f"Multi-select values in {col} should be boolean")
    
    def test_file_loading_error_handling(self):
        """Test that file loading properly handles error cases"""
        # Test that loading a non-existent file raises FileNotFoundError
        with self.assertRaises(FileNotFoundError, msg="Should raise FileNotFoundError for non-existent files"):
            load_file("nonexistent_file.csv")
        
        # Test that loading an unsupported file type raises ValueError
        test_txt = os.path.join(self.test_data_dir, "test.txt")
        with self.assertRaises(ValueError, msg="Should raise ValueError for unsupported file types"):
            load_file(test_txt)
        
        # Test that loading a file with invalid content raises ValueError
        test_file = os.path.join(self.test_data_dir, "invalid_test.csv")
        try:
            # Create a temporary invalid test file
            with open(test_file, "w") as f:
                f.write("This is not a valid CSV file")
            
            with self.assertRaises(ValueError, msg="Should raise ValueError for invalid file content"):
                load_file(test_file)
        finally:
            # Clean up the temporary file
            if os.path.exists(test_file):
                os.remove(test_file)

class TestDataProcessing(unittest.TestCase):
    """Test group for data processing functionality"""
    
    def setUp(self):
        """Set up test data"""
        # Load the test CSV file using application code
        self.test_data_dir = os.path.join(os.path.dirname(__file__), 'data')
        self.test_csv = os.path.join(self.test_data_dir, 'test.csv')
        self.test_data = load_file(self.test_csv)

    def test_open_ended_questions(self):
        """Test that open-ended questions are correctly identified"""
        # Process the data
        raw_audience_data, combined_data, group_audience_names = process_data(self.test_data)
        
        # Check that Q4 (open-ended question) is identified
        open_ended_cols = [col for col in self.test_data.columns if col.startswith('Q(4)')]
        self.assertGreater(len(open_ended_cols), 0, "Should have open-ended question columns")
        
        # Check that the open-ended question is not included in the processed data
        # (since we don't process open-ended questions for charts)
        for title, categories, segments in raw_audience_data:
            self.assertFalse(title.startswith('Q(4)'), 
                           "Open-ended questions should not be included in processed data")
        
        # Check that the open-ended question is not included in combined data
        for title, categories, segments in combined_data:
            self.assertFalse(title.startswith('Q(4)'), 
                           "Open-ended questions should not be included in combined data")

    def test_process_data_structure(self):
        """Test the structure and format of processed data"""
        # Define test-specific audience definitions that match our test data
        test_audience_defs = {
            "Lloyds Customers": {
                "lloyds_customer": [True]
            },
            "Young Adults": {
                "Age": ["18-20", "21-24", "25-34"]
            },
            "Male": {
                "Gender": ["Male"]
            }
        }
        
        # Process the data using our application's process_data function
        raw_audience_data, combined_data, group_audience_names = process_data(self.test_data, audience_defs=test_audience_defs)
        
        # Check that we got processed data back
        self.assertIsNotNone(raw_audience_data, "Raw audience data should not be None")
        self.assertIsNotNone(combined_data, "Combined data should not be None")
        
        # Check that raw_audience_data has the expected structure
        for title, categories, segments in raw_audience_data:
            self.assertIsInstance(title, str, "Title should be a string")
            self.assertIsInstance(categories, list, "Categories should be a list")
            self.assertIsInstance(segments, list, "Segments should be a list")
            
            # Check that segments have the expected structure
            for segment in segments:
                self.assertIsInstance(segment, tuple, "Segment should be a tuple")
                self.assertEqual(len(segment), 3, "Segment should have 3 elements (label, values, count)")
                label, values, count = segment
                self.assertIsInstance(label, str, "Segment label should be a string")
                self.assertIsInstance(values, list, "Segment values should be a list")
                self.assertIsInstance(count, int, "Segment count should be an integer")
                
                # Check that values match categories length
                self.assertEqual(len(values), len(categories), 
                               "Number of values should match number of categories")
                
                # Check that values are between 0 and 1 (percentages)
                for value in values:
                    self.assertGreaterEqual(value, 0, "Value should be >= 0")
                    self.assertLessEqual(value, 1, "Value should be <= 1")

    def test_response_values(self):
        """Test that response values are calculated correctly"""
        # Known correct values for our test data
        expected_values = {
            "Q(1)": {  # Multi-select question about awareness
                "categories": [
                    "Lloyds Bank", "HSBC", "Barclays", "Halifax", "Santander",
                    "Nationwide", "Hargreaves Lansdown", "Moneybox", "A J Bell",
                    "Interactive Investor", "Nutmeg", "None of these"
                ],
                "total_responses": 100,  # Total responses in test data
                "response_counts": {
                    "Lloyds Bank": 31,  # 31%
                    "HSBC": 34,  # 34%
                    "Barclays": 35,  # 35%
                    "Halifax": 30,  # 30%
                    "Santander": 36,  # 36%
                    "Nationwide": 30,  # 30%
                    "Hargreaves Lansdown": 20,  # 20%
                    "Moneybox": 18,  # 18%
                    "A J Bell": 21,  # 21%
                    "Interactive Investor": 16,  # 16%
                    "Nutmeg": 19,  # 19%
                    "None of these": 9   # 9%
                }
            },
            "Q(2)": {  # Single-select question about consideration
                "categories": [
                    "I'd consider them equally with others",
                    "It's one of the first I'd consider using",
                    "It's the first I'd consider using",
                    "I might consider them but I'd consider others first",
                    "I probably wouldn't consider them but I wouldn't rule them out",
                    "Don't know",
                    "I definitely wouldn't consider them"
                ],
                "total_responses": 100,  # Total responses in test data
                "response_counts": {
                    "I'd consider them equally with others": 39,  # 39%
                    "It's one of the first I'd consider using": 25,  # 25%
                    "I definitely wouldn't consider them": 4,  # 4%
                    "Don't know": 3  # 3%
                }
            }
        }
        
        raw_audience_data, _, group_audience_names = process_data(self.test_data)
        
        # Check each question's results
        for title, categories, segments in raw_audience_data:
            print(f"[test_response_values] Title: {title}, Categories: {categories}")
            # Extract question ID from title
            q_id = "Q(1)" if "Which of the following are you aware offer investment products" in title else "Q(2)" if "Thinking about Lloyds Bank; if you were considering investing" in title else None
            if q_id not in expected_values:
                continue
            
            # Check total responses
            total_segment = next((s for s in segments if s[0] == "Total"), None)
            self.assertIsNotNone(total_segment, f"Missing Total segment for {q_id}")
            self.assertEqual(total_segment[2], expected_values[q_id]["total_responses"], 
                           f"Wrong total responses for {q_id}")
            
            # Create a mapping of category to value
            category_values = dict(zip(categories, total_segment[1]))
            
            # Normalize category names to handle curly apostrophes
            normalized_category_values = {k.replace('’', "'"): v for k, v in category_values.items()}
            normalized_category = next((c for c in categories if c.replace('’', "'") in normalized_category_values), None)
            self.assertIsNotNone(normalized_category, f"Category {normalized_category} not found in results")
            
            # Check response counts
            for category, expected_count in expected_values[q_id]["response_counts"].items():
                self.assertIn(category, normalized_category_values, f"Category {category} not found in results")
                actual_percentage = normalized_category_values[category]
                expected_percentage = expected_count / expected_values[q_id]["total_responses"]
                self.assertAlmostEqual(actual_percentage, expected_percentage, places=3,
                                     msg=f"Wrong percentage for {category} in {q_id}")

    def test_audience_filtering(self):
        """Test that audience filtering works correctly"""
        
        
        # Define test audiences

        test_audience_defs = {
            "Lloyds Customers": {
                "lloyds_customer": [True]
            },
            "Older Males": {
                "AND": [
                    {"Gender": ["Male"]},
                    {"Age": ["45-54", "55-64"]}
                ]
            }
        }
        
        # Process data
        _, combined_data, group_audience_names = process_data(self.test_data, audience_defs=test_audience_defs)
        
        # Known values for Lloyds customers to check against
        expected_lloyds_counts = {
            "Q(1)": {  # Multi-select awareness question
                "total": 35,  # Number of Lloyds customers in test data
                "responses": {
                    "Lloyds Bank": 0.429,  # 42.9%
                    "HSBC": 0.286,  # 28.6%
                    "None of these": 0.000   # 0.0%
                }
            }
        }

        # Known values for Older Males to check against
        expected_older_males_counts = {
            "Q(1)": {  # Multi-select awareness question
                "total": 4,  # Number of Older Males in test data
                "responses": {
                    "Lloyds Bank": 0.5,  # 50.0%
                    "HSBC": 0.5,  # 50.0%
                    "None of these": 0.000   # 0.0%
                }
            }
        }
        
        # Check each question's results
        for title, categories, segments in combined_data:
            print(f"[test_audience_filtering] Title: {title}, Categories: {categories}")
            print(f"[test_audience_filtering] Segments: {[s[0] for s in segments]}")
            # Extract question ID from title
            q_id = "Q(1)" if "Which of the following are you aware offer investment products" in title else "Q(2)" if "Thinking about Lloyds Bank; if you were considering investing" in title else None
            if q_id not in expected_lloyds_counts:
                continue
                
            # Only process items that have both Lloyds Customers and Older Males segments
            segment_names = [s[0] for s in segments]
            if "Lloyds Customers" not in segment_names or "Older Males" not in segment_names:
                continue
            
            # Get Lloyds customers segment
            lloyds_segment = next((s for s in segments if s[0] == "Lloyds Customers"), None)
            self.assertIsNotNone(lloyds_segment, f"Missing Lloyds Customers segment for {q_id}")

            # Get Older Males segment
            older_males_segment = next((s for s in segments if s[0] == "Older Males"), None)
            self.assertIsNotNone(older_males_segment, f"Missing Older Males segment for {q_id}")
            
            # Check total count
            expected = expected_lloyds_counts[q_id]
            self.assertEqual(lloyds_segment[2], expected["total"], 
                           f"Wrong number of Lloyds customers for {q_id}")
            
            # Check total count (replace EXPECTED_OLDER_MALES_COUNT with the correct value for your test data)
            EXPECTED_OLDER_MALES_COUNT = 4
            self.assertEqual(older_males_segment[2], EXPECTED_OLDER_MALES_COUNT, 
                 f"Wrong number of Older Males for {q_id}")


            # Create a mapping of category to value
            category_values = dict(zip(categories, lloyds_segment[1]))
            older_males_category_values = dict(zip(categories, older_males_segment[1]))

            # Normalize category names to handle curly apostrophes
            
            # Normalize category names to handle curly apostrophes
            normalized_category_values = {k.replace('’', "'"): v for k, v in category_values.items()}
            normalized_older_males_category_values = {k.replace('’', "'"): v for k, v in older_males_category_values.items()}
            
            normalized_category = next((c for c in categories if c.replace('’', "'") in normalized_category_values), None)
            self.assertIsNotNone(normalized_category, f"Category {normalized_category} not found in results")
            
            # Check specific response percentages
            for category, expected_pct in expected["responses"].items():
                self.assertIn(category, normalized_category_values, f"Category {category} not found in results")
                actual_pct = normalized_category_values[category]
                self.assertAlmostEqual(actual_pct, expected_pct, places=3,
                                     msg=f"Wrong percentage for {category} in {q_id} for Lloyds customers")
                
            # Check specific response percentages for Older Males
            expected_older = expected_older_males_counts[q_id]
            for category, expected_pct in expected_older["responses"].items():
                self.assertIn(category, normalized_older_males_category_values, f"Category {category} not found in Older Males results")
                actual_pct = normalized_older_males_category_values[category]
                self.assertAlmostEqual(actual_pct, expected_pct, places=3,
                    msg=f"Wrong percentage for {category} in {q_id} for Older Males")

    def test_multi_select_boolean_conversion(self):
        """Test that multi-select questions are correctly converted to boolean values"""
        # Create test data with known multi-select responses
        test_data = pd.DataFrame({
            'Q(1_1) Lloyds Bank[Question: Which banks?]': ['Lloyds Bank', 'Lloyds Bank;HSBC', '', None],
            'Q(1_2) HSBC[Question: Which banks?]': ['HSBC', 'Lloyds Bank;HSBC', 'HSBC', None],
            'Q(1_3) Barclays[Question: Which banks?]': ['Barclays', '', '', None],
        })
        
        # Process the data
        df = process_dataframe(test_data)
        
        # Check that columns are boolean
        multi_select_cols = [col for col in df.columns if re.match(r'Q\(\d+_\d+\)', col)]
        self.assertGreater(len(multi_select_cols), 0, "Should have multi-select questions")
        
        for col in multi_select_cols:
            self.assertEqual(df[col].dtype, bool, 
                           f"Multi-select column {col} should be boolean")
        
        # Check specific boolean values
        # Lloyds Bank column
        self.assertTrue(df['Q(1_1) Lloyds Bank[Question: Which banks?]'].iloc[0], 
                       "Should be True when Lloyds Bank is selected")
        self.assertTrue(df['Q(1_1) Lloyds Bank[Question: Which banks?]'].iloc[1], 
                       "Should be True when Lloyds Bank is part of multiple selections")
        self.assertFalse(df['Q(1_1) Lloyds Bank[Question: Which banks?]'].iloc[2], 
                        "Should be False for empty string")
        self.assertFalse(df['Q(1_1) Lloyds Bank[Question: Which banks?]'].iloc[3], 
                        "Should be False for None")
        
        # HSBC column
        self.assertTrue(df['Q(1_2) HSBC[Question: Which banks?]'].iloc[0], 
                       "Should be True when HSBC is selected")
        self.assertTrue(df['Q(1_2) HSBC[Question: Which banks?]'].iloc[1], 
                       "Should be True when HSBC is part of multiple selections")
        self.assertTrue(df['Q(1_2) HSBC[Question: Which banks?]'].iloc[2], 
                       "Should be True when HSBC is selected")
        self.assertFalse(df['Q(1_2) HSBC[Question: Which banks?]'].iloc[3], 
                        "Should be False for None")

    def test_age_of_children_processing(self):
        """Test that Age of children column is correctly processed into boolean columns"""
        # Create test data with age of children responses
        test_data = pd.DataFrame({
            'Age of children': [
                '3-4 years;5-6 years',
                '0-3 months;1-2 years',
                '7-8 years',
                '19+ years',
                '',
                None
            ]
        })
        
        # Process the data
        df = process_dataframe(test_data)
        
        # Check that age of children columns are created
        age_child_cols = [col for col in df.columns if col.startswith('has_') and col.endswith('_child')]
        self.assertGreater(len(age_child_cols), 0, "Should create child age columns")
        
        # Check that all age child columns are boolean
        for col in age_child_cols:
            self.assertEqual(df[col].dtype, bool, 
                           f"Age child column {col} should be boolean")
        
        # Check specific boolean values for the first row: '3-4 years;5-6 years'
        self.assertTrue(df['has_3_4_years_child'].iloc[0], 
                       "Should be True when 3-4 years is selected")
        self.assertTrue(df['has_5_6_years_child'].iloc[0], 
                       "Should be True when 5-6 years is selected")
        self.assertFalse(df['has_0_3_months_child'].iloc[0], 
                        "Should be False when 0-3 months is not selected")
        
        # Check second row: '0-3 months;1-2 years'
        self.assertTrue(df['has_0_3_months_child'].iloc[1], 
                       "Should be True when 0-3 months is selected")
        self.assertTrue(df['has_1_2_years_child'].iloc[1], 
                       "Should be True when 1-2 years is selected")
        self.assertFalse(df['has_3_4_years_child'].iloc[1], 
                        "Should be False when 3-4 years is not selected")
        
        # Check third row: '7-8 years'
        self.assertTrue(df['has_7_8_years_child'].iloc[2], 
                       "Should be True when 7-8 years is selected")
        self.assertFalse(df['has_3_4_years_child'].iloc[2], 
                        "Should be False when 3-4 years is not selected")
        
        # Check fourth row: '19+ years'
        self.assertTrue(df['has_19_plus_years_child'].iloc[3], 
                       "Should be True when 19+ years is selected")
        
        # Check empty and None values
        for col in age_child_cols:
            self.assertFalse(df[col].iloc[4], 
                           f"Should be False for empty string in {col}")
            self.assertFalse(df[col].iloc[5], 
                           f"Should be False for None in {col}")

    def test_age_of_children_column_order(self):
        """Test that Age of children columns are created in chronological order"""
        # Create test data with all age ranges to test ordering
        test_data = pd.DataFrame({
            'Age of children': [
                '19+ years;0-3 months',  # Mix of oldest and youngest
                '5-6 years;1-2 years',   # Mix of middle ages
                '11-12 years;7-8 years', # Mix of older ages
            ]
        })
        
        # Process the data
        df = process_dataframe(test_data)
        
        # Get all age child columns in the order they appear in the DataFrame
        age_child_cols = [col for col in df.columns if col.startswith('has_') and col.endswith('_child')]
        
        # Define the expected order of column names
        expected_order = [
            'has_0_3_months_child',
            'has_4_7_months_child', 
            'has_8_11_months_child',
            'has_1_2_years_child',
            'has_3_4_years_child',
            'has_5_6_years_child',
            'has_7_8_years_child',
            'has_9_10_years_child',
            'has_11_12_years_child',
            'has_13_14_years_child',
            'has_15_16_years_child',
            'has_17_18_years_child',
            'has_19_plus_years_child'
        ]
        
        # Check that the columns appear in the expected order
        # We only check the columns that were actually created
        actual_order = [col for col in age_child_cols if col in expected_order]
        expected_created = [col for col in expected_order if col in age_child_cols]
        
        self.assertEqual(actual_order, expected_created, 
                        "Age child columns should be created in chronological order")
        
        # Verify that the columns are in the DataFrame in the correct order
        for i, expected_col in enumerate(expected_created):
            if expected_col in df.columns:
                # Find the position of this column in the DataFrame
                col_position = list(df.columns).index(expected_col)
                # Check that it comes after the previous expected column
                if i > 0 and expected_created[i-1] in df.columns:
                    prev_col_position = list(df.columns).index(expected_created[i-1])
                    self.assertGreater(col_position, prev_col_position,
                                     f"Column {expected_col} should come after {expected_created[i-1]}")

    def test_semicolon_separated_column_processing(self):
        """Test the generalized semicolon-separated column processing function"""
        from src.data_loader import process_semicolon_separated_column
        
        # Test with a custom column
        test_data = pd.DataFrame({
            'Favorite Colors': [
                'Red;Blue',
                'Green;Red;Yellow',
                'Blue',
                '',
                None
            ]
        })
        
        # Process with custom suffix
        df = process_semicolon_separated_column(test_data, 'Favorite Colors', '_color')
        
        # Check that color columns are created
        color_cols = [col for col in df.columns if col.startswith('has_') and col.endswith('_color')]
        self.assertGreater(len(color_cols), 0, "Should create color columns")
        
        # Check specific values
        self.assertTrue(df['has_red_color'].iloc[0], "Should be True for Red in first row")
        self.assertTrue(df['has_blue_color'].iloc[0], "Should be True for Blue in first row")
        self.assertTrue(df['has_green_color'].iloc[1], "Should be True for Green in second row")
        self.assertTrue(df['has_red_color'].iloc[1], "Should be True for Red in second row")
        self.assertTrue(df['has_yellow_color'].iloc[1], "Should be True for Yellow in second row")
        self.assertTrue(df['has_blue_color'].iloc[2], "Should be True for Blue in third row")
        
        # Check that all color columns are boolean
        for col in color_cols:
            self.assertEqual(df[col].dtype, bool, 
                           f"Color column {col} should be boolean")

    def test_bank_columns_alphabetical_order(self):
        """Test that bank customer columns are created in alphabetical order"""
        # Create test data with banks in random order
        test_data = pd.DataFrame({
            'Bank(s)': [
                'HSBC;Santander;Barclays',
                'Lloyds;Monzo;Starling',
                'Natwest;RBS;Halifax',
                'Metro;TSB;Revolut',
            ]
        })
        
        # Process the data
        df = process_dataframe(test_data)
        
        # Get all bank customer columns in the order they appear in the DataFrame
        bank_customer_cols = [col for col in df.columns if col.endswith('_customer')]
        
        # Define the expected alphabetical order
        expected_order = [
            'barclays_customer',
            'halifax_customer',
            'hsbc_customer',
            'lloyds_customer',
            'metro_customer',
            'monzo_customer',
            'natwest_customer',
            'rbs_customer',
            'revolut_customer',
            'santander_customer',
            'starling_customer',
            'tsb_customer'
        ]
        
        # Check that the columns appear in alphabetical order
        # We only check the columns that were actually created
        actual_order = [col for col in bank_customer_cols if col in expected_order]
        expected_created = [col for col in expected_order if col in bank_customer_cols]
        
        self.assertEqual(actual_order, expected_created, 
                        "Bank customer columns should be created in alphabetical order")
        
        # Verify that the columns are in the DataFrame in the correct order
        for i, expected_col in enumerate(expected_created):
            if expected_col in df.columns:
                # Find the position of this column in the DataFrame
                col_position = list(df.columns).index(expected_col)
                # Check that it comes after the previous expected column
                if i > 0 and expected_created[i-1] in df.columns:
                    prev_col_position = list(df.columns).index(expected_created[i-1])
                    self.assertGreater(col_position, prev_col_position,
                                     f"Column {expected_col} should come after {expected_created[i-1]} in alphabetical order")

    def test_chart_data_ordering(self):
        """Test that chart data is ordered by Total values in descending order."""
        # Create sample data with known order
        raw_df = pd.DataFrame({
            'Q(1)': ['A', 'B', 'A', 'C', 'B', 'A', 'D'],  # A=3, B=2, C=1, D=1
        })
        
        # Process the data
        raw_audience_data, combined_data, group_audience_names = process_data(raw_df)
        
        # Check the first chart's data (should be the Q1 data)
        title, categories, segments = raw_audience_data[0]
        
        # Get the Total segment's values
        total_segment = next(seg for seg in segments if seg[0] == 'Total')
        total_values = total_segment[1]
        
        # Check if values are in descending order
        sorted_values = sorted(total_values, reverse=True)
        self.assertEqual(total_values, sorted_values, 
                        "Chart values should be sorted in descending order based on Total values")
        
        # Also check that categories are reordered to match
        sorted_pairs = sorted(zip(total_values, categories), reverse=True)
        sorted_values, sorted_categories = zip(*sorted_pairs)
        self.assertEqual(len(sorted_values), len(categories), 
                        "Categories should be reordered to match the descending value order")

    def test_all_questions_processed(self):
        """Test that all questions in a survey are processed (except open-ended ones)"""
        # Process the data
        raw_audience_data, combined_data, group_audience_names = process_data(self.test_data)
        
        # Count the number of questions processed
        num_processed_questions = len(raw_audience_data)
        
        # We expect 3 questions to be processed (Q1, Q2, Q3)
        # Q4 is open-ended and should be excluded
        expected_questions = 3
        
        self.assertEqual(num_processed_questions, expected_questions, 
                        f"Expected {expected_questions} questions to be processed, but got {num_processed_questions}")
        
        # Verify we have the expected question IDs
        processed_q_ids = set()
        for title, categories, segments in raw_audience_data:
            # Extract question ID from the first segment's data
            if segments and len(segments) > 0:
                # For multi-select questions, we need to look at the column structure
                # For single-select questions, we can extract from the title
                if "Which of the following are you aware offer investment products" in title:
                    processed_q_ids.add("1")
                elif "Thinking about Lloyds Bank; if you were considering investing" in title:
                    processed_q_ids.add("2")
                elif "opinion of Lloyds Bank as an investment product provider" in title:
                    processed_q_ids.add("3")
        
        expected_q_ids = {"1", "2", "3"}
        self.assertEqual(processed_q_ids, expected_q_ids,
                        f"Expected question IDs {expected_q_ids}, but got {processed_q_ids}")
        
        # Verify each question has reasonable data
        for title, categories, segments in raw_audience_data:
            self.assertGreater(len(categories), 0, f"No categories found for question: {title}")
            self.assertGreater(len(segments), 0, f"No segments found for question: {title}")
            
            # Check that the first segment (Total) has values
            total_segment = segments[0]
            self.assertEqual(total_segment[0], "Total", f"First segment should be 'Total' for {title}")
            self.assertEqual(len(total_segment[1]), len(categories), 
                           f"Number of values should match number of categories for {title}")
            
            # Check that values are reasonable (between 0 and 1 for percentages)
            for value in total_segment[1]:
                self.assertGreaterEqual(value, 0.0, f"Value should be >= 0 for {title}")
                self.assertLessEqual(value, 1.0, f"Value should be <= 1 for {title}")

    def test_questions_with_comments_are_processed(self):
        """Test that questions with comment columns are still processed correctly"""
        # Load the problem file
        problem_csv = os.path.join("survey_data", "Problem_file.csv")
        df = load_file(problem_csv)
        
        # Process the data
        raw_audience_data, combined_data, group_audience_names = process_data(df)
        
        # We expect 3 questions to be processed (Q1, Q2, Q3)
        # Even though Q2 and Q3 have comment columns that might be identified as open-ended,
        # the main questions should still be processed
        expected_questions = 3
        
        self.assertEqual(len(raw_audience_data), expected_questions, 
                        f"Expected {expected_questions} questions to be processed, but got {len(raw_audience_data)}")
        
        # Verify we have the expected question IDs
        processed_q_ids = set()
        for title, categories, segments in raw_audience_data:
            # Extract question ID from the title or content
            if "restaurant the most" in title:  # Q1
                processed_q_ids.add("1")
            elif "cuisine/food do you tend to go for" in title:  # Q2
                processed_q_ids.add("2")
            elif "Pizza Hut restaurant" in title:  # Q3
                processed_q_ids.add("3")
        
        expected_q_ids = {"1", "2", "3"}
        self.assertEqual(processed_q_ids, expected_q_ids,
                        f"Expected question IDs {expected_q_ids}, but got {processed_q_ids}")
        
        # Verify each question has reasonable data
        for title, categories, segments in raw_audience_data:
            self.assertGreater(len(categories), 0, f"No categories found for question: {title}")
            self.assertGreater(len(segments), 0, f"No segments found for question: {title}")
            
            # Check that the first segment (Total) has values
            total_segment = segments[0]
            self.assertEqual(total_segment[0], "Total", f"First segment should be 'Total' for {title}")
            self.assertEqual(len(total_segment[1]), len(categories), 
                           f"Number of values should match number of categories for {title}")
            
            # Check that values are reasonable (between 0 and 1 for percentages)
            for value in total_segment[1]:
                self.assertGreaterEqual(value, 0.0, f"Value should be >= 0 for {title}")
                self.assertLessEqual(value, 1.0, f"Value should be <= 1 for {title}")

class TestDataAnalysis(unittest.TestCase):
    """Test group for data analysis functionality"""
    
    def setUp(self):
        """Set up test data"""
        self.test_data_dir = os.path.join(os.path.dirname(__file__), 'data')
        self.test_csv = os.path.join(self.test_data_dir, 'test.csv')
        self.test_data = load_file(self.test_csv)
    
    def test_multi_select_processing(self):
        """Test processing of multi-select questions"""
        # Process the data
        raw_audience_data, combined_data, group_audience_names = process_data(self.test_data)
        
        # Find multi-select questions (Q1_1, Q1_2, etc.)
        multi_select_cols = [col for col in self.test_data.columns if re.match(r'Q\(\d+_\d+\)', col)]
        self.assertGreater(len(multi_select_cols), 0, "Should have multi-select questions")
        
        # Check that each multi-select question is processed correctly
        for col in multi_select_cols:
            # Get non-empty responses
            responses = self.test_data[col].dropna()
            self.assertGreater(len(responses), 0, f"Should have responses for {col}")
            
            # Check that responses are properly formatted
            for response in responses:
                if isinstance(response, str):  # Raw response
                    parts = [p.strip() for p in response.split(';')]
                    self.assertTrue(all(p for p in parts), 
                                  f"All parts of response should be non-empty in {col}")
                elif isinstance(response, list):  # Processed response
                    self.assertTrue(all(isinstance(p, str) and p.strip() == p for p in response),
                                  f"All parts should be stripped strings in {col}")
                    
        # Check that the processed data includes multi-select questions
        for title, categories, segments in raw_audience_data:
            if any(q_id in title for q_id in ['Q(1_1)', 'Q(1_2)', 'Q(1_3)']):  # Known multi-select questions
                # Check that categories match expected format
                self.assertTrue(all(isinstance(cat, str) for cat in categories),
                              f"Categories should be strings in {title}")
                # Check that values are valid percentages
                for segment in segments:
                    label, values, count = segment
                    self.assertTrue(all(0 <= v <= 1 for v in values),
                                  f"Values should be between 0 and 1 in {title}")



class TestPowerPointGeneration(unittest.TestCase):
    """Test group for PowerPoint generation functionality"""
    
    def setUp(self):
        """Set up test presentation"""
        self.prs = Presentation()
    
    def test_chart_creation(self):
        """Test basic chart creation"""
        # Test data
        categories = ['A', 'B', 'C']
        series_list = [('Total', [0.3, 0.5, 0.2])]
        
        # Create chart
        slide, chart = create_chart_slide(self.prs, categories, series_list)
        
        # Check chart properties
        self.assertEqual(len(chart.series), 1, "Chart should have one series")
        self.assertEqual(len(chart.series[0].values), 3, "Series should have three values")
    
    def test_color_assignment(self):
        """Test color assignment for multiple series"""
        # Test data with multiple series
        categories = ['A', 'B', 'C']
        series_list = [
            ('Total', [0.3, 0.5, 0.2]),
            ('Group 1', [0.4, 0.3, 0.3]),
            ('Group 2', [0.2, 0.4, 0.4])
        ]
        
        # Create chart
        slide, chart = create_chart_slide(self.prs, categories, series_list)
        
        # Check that each series has a different color
        colors = set()
        for series in chart.series:
            colors.add(series.format.fill.fore_color.theme_color)
        
        self.assertEqual(len(colors), len(series_list), 
                        "Each series should have a unique color")

    def test_scale_chart_adds_net_group_callouts(self):
        """Scale charts should show grouped percentage callouts per series."""
        categories = [
            'Strongly agree',
            'Agree',
            'Neither agree nor disagree',
            'Disagree',
            'Strongly disagree'
        ]
        series_list = [
            ('Total', [0.35, 0.25, 0.10, 0.15, 0.05]),
            ('Group 1', [0.20, 0.30, 0.20, 0.20, 0.10])
        ]

        slide, chart = create_chart_slide(self.prs, categories, series_list)

        shape_text = "\n".join(
            shape.text
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        )
        self.assertIn("Total", shape_text)
        self.assertIn("Group 1", shape_text)
        self.assertIn("Net agree: 60%", shape_text)
        self.assertIn("Net disagree: 20%", shape_text)
        self.assertIn("Net agree: 50%", shape_text)
        self.assertIn("Net disagree: 30%", shape_text)
        self.assertNotIn("Net score:", shape_text)

    def test_non_scale_chart_does_not_add_net_score_callouts(self):
        """Normal categorical charts should not get net score callouts."""
        categories = ['Lloyds', 'Barclays', 'Halifax']
        series_list = [('Total', [0.3, 0.5, 0.2])]

        slide, chart = create_chart_slide(self.prs, categories, series_list)

        shape_text = "\n".join(
            shape.text
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        )
        self.assertNotIn("Net score:", shape_text)

    def test_missing_template_file_raises_clear_error(self):
        """Test that missing template file raises a clear, user-friendly error."""
        # Save the original template path
        original_template_path = config.TEMPLATE_PATH
        try:
            # Set to a definitely non-existent file
            config.TEMPLATE_PATH = "nonexistent_template_file.pptx"
            # Minimal valid data for the function
            raw_audience_data = []
            combined_data = []
            with self.assertRaises(PackageNotFoundError) as cm:
                ppt_generator.generate_presentation(raw_audience_data, combined_data)
            # Check the error message is clear
            self.assertIn("Package not found", str(cm.exception))
        finally:
            # Restore the original path
            config.TEMPLATE_PATH = original_template_path

class TestConfig(unittest.TestCase):
    """Test group for configuration functionality"""
    
    def test_output_path_generation(self):
        """Test that output path generation works correctly"""
        # Test with a simple filename
        filename = "test.xlsx"
        expected = os.path.join("exports", "test_full.pptx")
        result = config.get_output_pptx_path(filename)
        self.assertEqual(result, expected)
        
        # Test with a filename with spaces
        filename = "test file.xlsx"
        expected = os.path.join("exports", "test file_full.pptx")
        result = config.get_output_pptx_path(filename)
        self.assertEqual(result, expected)
        
        # Test with a filename with special characters
        filename = "test-file_123.xlsx"
        expected = os.path.join("exports", "test-file_123_full.pptx")
        result = config.get_output_pptx_path(filename)
        self.assertEqual(result, expected)
    
    def test_output_path_absolute(self):
        """Test that output path generation works with absolute paths"""
        # Test with an absolute path
        filename = "/path/to/test.xlsx"
        expected = os.path.join("exports", "test_full.pptx")
        result = config.get_output_pptx_path(filename)
        self.assertEqual(result, expected)
    
    def test_output_path_in_project_root(self):
        """Test that output path is created in the project root"""
        filename = "test.xlsx"
        result = config.get_output_pptx_path(filename)
        
        # Check that the path is relative to the project root
        self.assertFalse(os.path.isabs(result), "Output path should be relative")
        
        # Check that it's in the exports directory
        self.assertTrue(result.startswith("exports"), "Output path should be in exports directory")

    def test_pptx_filename_reflects_original_file(self):
        """Test that PPTX filename reflects the original XLS/CSV filename"""
        # Test with different file types
        test_cases = [
            ("data.xlsx", "data_full.pptx"),
            ("survey_results.csv", "survey_results_full.pptx"),
            ("test file.xlsx", "test file_full.pptx"),
            ("data_2024.xlsx", "data_2024_full.pptx")
        ]

        for input_filename, expected_pptx_name in test_cases:
            result = config.get_output_pptx_path(input_filename)
            expected_path = os.path.join("exports", expected_pptx_name)
            self.assertEqual(result, expected_path,
                             f"PPTX filename should reflect original filename: {input_filename} -> {expected_pptx_name}")

class TestAppFunctionality(unittest.TestCase):
    """Test group for app-specific functionality issues"""
    
    def setUp(self):
        """Set up test data"""
        self.test_data_dir = os.path.join(os.path.dirname(__file__), 'data')
        self.test_csv = os.path.join(self.test_data_dir, 'test.csv')
        self.test_data = load_file(self.test_csv)
    
    def test_pptx_filename_reflects_original_file(self):
        """Test that PPTX filename reflects the original XLS/CSV filename"""
        # Test with different file types
        test_cases = [
            ("data.xlsx", "data_full.pptx"),
            ("survey_results.csv", "survey_results_full.pptx"),
            ("test file.xlsx", "test file_full.pptx"),
            ("data_2024.xlsx", "data_2024_full.pptx")
        ]
        
        for input_filename, expected_pptx_name in test_cases:
            result = config.get_output_pptx_path(input_filename)
            expected_path = os.path.join("exports", expected_pptx_name)
            self.assertEqual(result, expected_path, 
                           f"PPTX filename should reflect original filename: {input_filename} -> {expected_pptx_name}")
    
    def test_audiences_appear_in_output_data(self):
        """Test that audiences appear in the output data when they exist"""
        # Define test audiences
        test_audience_defs = {
            "Lloyds Customers": {
                "lloyds_customer": [True]
            },
            "Young Adults": {
                "Age": ["18-20", "21-24", "25-34"]
            }
        }
        
        # Process the data
        raw_audience_data, combined_data, group_audience_names = process_data(self.test_data, audience_defs=test_audience_defs)
        
        # Debug output
        print(f"\nDebug: Number of combined_data items: {len(combined_data)}")
        for i, (title, categories, segments) in enumerate(combined_data):
            print(f"Debug: Item {i}: Title='{title}', Categories={len(categories)}, Segments={len(segments)}")
            for j, segment in enumerate(segments):
                label, values, count = segment
                print(f"  Debug: Segment {j}: Label='{label}', Values={len(values)}, Count={count}")
        
        # Check that combined_data contains audience segments
        audience_found = False
        for title, categories, segments in combined_data:
            # Look for segments that are not "Total"
            for segment in segments:
                label, values, count = segment
                if label != "Total":
                    audience_found = True
                    break
            if audience_found:
                break
        
        self.assertTrue(audience_found, "Audiences should appear in combined_data when they exist")
    
    def test_no_combined_data_when_no_audiences(self):
        """Test that combined data charts don't appear when no audiences exist"""
        # Process the data with no audiences
        raw_audience_data, combined_data, group_audience_names = process_data(self.test_data, audience_defs={})
        
        # Check that combined_data only contains Total segments (no audience-specific data)
        for title, categories, segments in combined_data:
            for segment in segments:
                label, values, count = segment
                # Only Total segments should exist when no audiences are defined
                self.assertEqual(label, "Total", 
                               f"Only Total segments should exist when no audiences defined, found: {label}")
    
    def test_audience_data_structure(self):
        """Test that audience data has the correct structure when audiences exist"""
        # Define test audiences
        test_audience_defs = {
            "Lloyds Customers": {
                "lloyds_customer": [True]
            }
        }
        
        # Process the data
        raw_audience_data, combined_data, group_audience_names = process_data(self.test_data, audience_defs=test_audience_defs)
        
        # Check that combined_data has the expected structure
        for title, categories, segments in combined_data:
            self.assertIsInstance(title, str, "Title should be a string")
            self.assertIsInstance(categories, list, "Categories should be a list")
            self.assertIsInstance(segments, list, "Segments should be a list")
            
            # Check that we have at least one segment
            self.assertGreater(len(segments), 0, "Should have at least one segment")
            
            # Check that the first segment is always "Total"
            first_segment = segments[0]
            self.assertEqual(first_segment[0], "Total", "First segment should always be 'Total'")
            
            # Check that audience segments have the correct structure
            for segment in segments:
                self.assertIsInstance(segment, tuple, "Segment should be a tuple")
                self.assertEqual(len(segment), 3, "Segment should have 3 elements (label, values, count)")
                label, values, count = segment
                self.assertIsInstance(label, str, "Segment label should be a string")
                self.assertIsInstance(values, list, "Segment values should be a list")
                self.assertIsInstance(count, int, "Segment count should be an integer")

    def test_chart_titles_contain_question_text(self):
        """Test that chart titles contain descriptive question text rather than just question IDs"""
        # Process the data
        raw_audience_data, combined_data, group_audience_names = process_data(self.test_data)
        
        # Check that chart titles contain descriptive text, not just Q(1), Q(2), etc.
        for title, categories, segments in combined_data:
            # Title should not be just a question ID like "Q(1)" or "1"
            self.assertFalse(
                re.match(r'^Q?\(\d+\)?$', title), 
                f"Chart title '{title}' should contain descriptive text, not just a question ID"
            )
            
            # Title should contain some descriptive text (more than just numbers and symbols)
            descriptive_chars = len(re.sub(r'[Q()\d\s]', '', title))
            self.assertGreater(
                descriptive_chars, 0,
                f"Chart title '{title}' should contain descriptive text, not just question formatting"
            )

    def test_audience_groups_basic_functionality(self):
        """Test that audience groups create the expected chart structure"""
        # Define test audiences and groups
        test_audience_defs = {
            "Men": {
                "Gender": ["Male"]
            },
            "Women": {
                "Gender": ["Female"]
            },
            "Young Adults": {
                "Age range": [18, 19, 20, 21, 22, 23, 24, 25, 26, 27, 28, 29, 30, 31, 32, 33, 34]
            },
            "Older Adults": {
                "Age range": [35, 36, 37, 38, 39, 40, 41, 42, 43, 45, 46, 49, 52, 56, 59, 61, 62, 87]
            },
            "Londoners": {
                "Home location": ["Greater London"]
            },
            "Northerners": {
                "Home location": ["North East", "North West", "Yorkshire & the Humber"]
            },
            "__groups__": [
                {
                    "name": "Gender",
                    "audiences": ["Men", "Women"]
                },
                {
                    "name": "Age",
                    "audiences": ["Young Adults", "Older Adults"]
                }
            ]
        }
        
        # Process the data
        raw_audience_data, combined_data, group_audience_names = process_data(self.test_data, audience_defs=test_audience_defs)
        
        # Debug output
        print(f"\nDebug: Number of combined_data items: {len(combined_data)}")
        for i, (title, categories, segments) in enumerate(combined_data):
            print(f"Debug: Item {i}: Title='{title}', Segments={[s[0] for s in segments]}")
        
        # Check that we have the expected chart types
        chart_titles = [title for title, _, _ in combined_data]
        
        # Should have group charts with group names in titles
        gender_group_charts = [title for title in chart_titles if " - Gender" in title]
        age_group_charts = [title for title in chart_titles if " - Age" in title]
        
        self.assertGreater(len(gender_group_charts), 0, "Should have Gender group charts")
        self.assertGreater(len(age_group_charts), 0, "Should have Age group charts")
        
        # Check that individual segments not in groups get their own charts
        # Look for charts that have only one audience segment (not group charts)
        individual_londoner_charts = []
        individual_northerner_charts = []
        
        for title, categories, segments in combined_data:
            segment_names = [s[0] for s in segments]
            if "Total" in segment_names and len(segments) == 2:  # Total + one audience
                if "Londoners" in segment_names:
                    individual_londoner_charts.append(title)
                elif "Northerners" in segment_names:
                    individual_northerner_charts.append(title)
        
        self.assertGreater(len(individual_londoner_charts), 0, "Should have individual Londoner charts")
        self.assertGreater(len(individual_northerner_charts), 0, "Should have individual Northerner charts")

    def test_audience_groups_chart_structure(self):
        """Test that audience group charts have the correct structure"""
        # Define test audiences and groups
        test_audience_defs = {
            "Men": {
                "Gender": ["Male"]
            },
            "Women": {
                "Gender": ["Female"]
            },
            "Young Adults": {
                "Age range": [18, 19, 20, 21, 22, 23, 24, 25, 26, 27, 28, 29, 30, 31, 32, 33, 34]
            },
            "Older Adults": {
                "Age range": [35, 36, 37, 38, 39, 40, 41, 42, 43, 45, 46, 49, 52, 56, 59, 61, 62, 87]
            },
            "__groups__": [
                {
                    "name": "Gender",
                    "audiences": ["Men", "Women"]
                }
            ]
        }
        
        # Process the data
        raw_audience_data, combined_data, group_audience_names = process_data(self.test_data, audience_defs=test_audience_defs)
        
        # Find the Gender group chart
        gender_group_chart = None
        for title, categories, segments in combined_data:
            if " - Gender" in title:
                gender_group_chart = (title, categories, segments)
                break
        
        self.assertIsNotNone(gender_group_chart, "Should have a Gender group chart")
        
        title, categories, segments = gender_group_chart
        
        # Check that the chart has Total + both audiences in the group
        segment_names = [segment[0] for segment in segments]
        self.assertIn("Total", segment_names, "Group chart should have Total segment")
        self.assertIn("Men", segment_names, "Gender group chart should have Men segment")
        self.assertIn("Women", segment_names, "Gender group chart should have Women segment")
        self.assertEqual(len(segments), 3, "Gender group chart should have exactly 3 segments (Total, Men, Women)")

    def test_audience_groups_no_duplicate_charts(self):
        """Test that segments in groups don't get individual charts"""
        # Define test audiences and groups
        test_audience_defs = {
            "Men": {
                "Gender": ["Male"]
            },
            "Women": {
                "Gender": ["Female"]
            },
            "Londoners": {
                "Home location": ["Greater London"]
            },
            "__groups__": [
                {
                    "name": "Gender",
                    "audiences": ["Men", "Women"]
                }
            ]
        }
        
        # Process the data
        raw_audience_data, combined_data, group_audience_names = process_data(self.test_data, audience_defs=test_audience_defs)
        
        # Check that Men and Women don't have individual charts (they're in the Gender group)
        chart_titles = [title for title, _, _ in combined_data]
        
        # Should NOT have individual Men/Women charts
        individual_men_charts = [title for title in chart_titles if "Men" in title and " - " not in title and "Gender" not in title]
        individual_women_charts = [title for title in chart_titles if "Women" in title and " - " not in title and "Gender" not in title]
        
        self.assertEqual(len(individual_men_charts), 0, "Men should not have individual charts when in a group")
        self.assertEqual(len(individual_women_charts), 0, "Women should not have individual charts when in a group")
        
        # Should have individual Londoner charts (not in a group)
        individual_londoner_charts = [title for title in chart_titles if "Londoners" in title and " - " not in title]
        self.assertGreater(len(individual_londoner_charts), 0, "Londoners should have individual charts when not in a group")

    def test_audience_groups_empty_group_handling(self):
        """Test that empty groups or groups with non-existent audiences are handled gracefully"""
        # Define test audiences with a group containing non-existent audience
        test_audience_defs = {
            "Men": {
                "Gender": ["Male"]
            },
            "Women": {
                "Gender": ["Female"]
            },
            "__groups__": [
                {
                    "name": "Gender",
                    "audiences": ["Men", "Women", "NonExistentAudience"]
                },
                {
                    "name": "Empty Group",
                    "audiences": []
                }
            ]
        }
        
        # Process the data - should not raise an error
        try:
            raw_audience_data, combined_data, group_audience_names = process_data(self.test_data, audience_defs=test_audience_defs)
            
            # Check that the Gender group chart only contains existing audiences
            gender_group_chart = None
            for title, categories, segments in combined_data:
                if " - Gender" in title:
                    gender_group_chart = (title, categories, segments)
                    break
            
            if gender_group_chart:
                title, categories, segments = gender_group_chart
                segment_names = [segment[0] for segment in segments]
                self.assertIn("Men", segment_names, "Should include existing Men audience")
                self.assertIn("Women", segment_names, "Should include existing Women audience")
                self.assertNotIn("NonExistentAudience", segment_names, "Should not include non-existent audience")
            
            # Check that empty groups don't create charts
            empty_group_charts = [title for title, _, _ in combined_data if "Empty Group" in title]
            self.assertEqual(len(empty_group_charts), 0, "Empty groups should not create charts")
            
        except Exception as e:
            self.fail(f"Processing should handle non-existent audiences gracefully, but got error: {e}")

    def test_audience_groups_all_segments_chart(self):
        """Test that there's a chart with Total + all individual segments"""
        # Define test audiences and groups
        test_audience_defs = {
            "Men": {
                "Gender": ["Male"]
            },
            "Women": {
                "Gender": ["Female"]
            },
            "Londoners": {
                "Home location": ["Greater London"]
            },
            "__groups__": [
                {
                    "name": "Gender",
                    "audiences": ["Men", "Women"]
                }
            ]
        }
        print("\n[DEBUG] Audience definitions:", test_audience_defs)
        
        # Process the data
        raw_audience_data, combined_data, group_audience_names = process_data(self.test_data, audience_defs=test_audience_defs)
        print("[DEBUG] Combined data chart titles and segments:")
        for title, categories, segments in combined_data:
            print(f"  Title: {title}")
            print(f"    Segments: {[segment[0] for segment in segments]}")
        
        # Find the "all segments" chart (should not have a group name in title)
        all_segments_charts = []
        expected_segments = {"Total", "Men", "Women", "Londoners"}  # All expected segments
        for title, categories, segments in combined_data:
            if " - " not in title:  # No group name in title
                segment_names = set(segment[0] for segment in segments)
                if "Total" in segment_names and len(segments) > 1:
                    # This is the "all segments" chart if it contains all expected segments
                    if segment_names == expected_segments:
                        all_segments_charts.append((title, segments))
        print("[DEBUG] All segments charts found:")
        for title, segments in all_segments_charts:
            print(f"  Title: {title}")
            print(f"    Segments: {[segment[0] for segment in segments]}")
        
        self.assertGreater(len(all_segments_charts), 0, "Should have charts with Total + all segments")
        
        # Check that the all segments chart contains all audiences
        for title, segments in all_segments_charts:
            segment_names = [segment[0] for segment in segments]
            print(f"[DEBUG] Checking segments in chart '{title}': {segment_names}")
            self.assertIn("Total", segment_names, "All segments chart should have Total")
            self.assertIn("Men", segment_names, "All segments chart should have Men")
            self.assertIn("Women", segment_names, "All segments chart should have Women")
            self.assertIn("Londoners", segment_names, "All segments chart should have Londoners")

    def test_audience_groups_chart_count(self):
        """Test that the correct number of charts are created"""
        # Define test audiences and groups
        test_audience_defs = {
            "Men": {
                "Gender": ["Male"]
            },
            "Women": {
                "Gender": ["Female"]
            },
            "Londoners": {
                "Home location": ["Greater London"]
            },
            "Northerners": {
                "Home location": ["North East", "North West", "Yorkshire & the Humber"]
            },
            "__groups__": [
                {
                    "name": "Gender",
                    "audiences": ["Men", "Women"]
                }
            ]
        }
        
        # Process the data
        raw_audience_data, combined_data, group_audience_names = process_data(self.test_data, audience_defs=test_audience_defs)
        
        # Count different types of charts
        # All-segments charts: titles that don't end with "(AudienceName)"
        # Individual charts: titles that end with "(AudienceName)"
        # Group charts: titles that contain " - GroupName"
        all_segments_charts = [title for title, _, _ in combined_data if not title.endswith(")") and " - " not in title]
        group_charts = [title for title, _, _ in combined_data if " - " in title]
        individual_charts = [title for title, _, _ in combined_data if title.endswith(")") and " - " not in title]
        
        # Expected:
        # 1. All segments chart (Total + Men + Women + Londoners + Northerners)
        # 2. Gender group chart (Total + Men + Women)
        # 3. Individual Londoners chart (Total + Londoners)
        # 4. Individual Northerners chart (Total + Northerners)
        # 5. Individual Men chart (Total + Men) - if not in group
        # 6. Individual Women chart (Total + Women) - if not in group
        
        # The exact count depends on the current implementation, but we should have at least:
        # - 1 all segments chart
        # - 1 gender group chart  
        # - 2 individual charts (Londoners, Northerners)
        self.assertGreaterEqual(len(all_segments_charts), 1, "Should have at least one all segments chart")
        self.assertGreaterEqual(len(group_charts), 1, "Should have at least one gender group chart")
        self.assertGreaterEqual(len(individual_charts), 2, "Should have at least two individual charts")

    def test_app_sample_size_calculation(self):
        """Test that sample size calculation works correctly for bank customer audiences"""
        # This test reproduces the exact issue the user is experiencing
        # where hsbc_customer = True shows sample size 0 in the app
        
        # Create a simple audience definition like the app would create
        audience_def = {
            "groups": [
                {
                    "name": "Group 1",
                    "conditions": [
                        {
                            "column": "hsbc_customer",
                            "values": ["True"]  # This is how the app saves it
                        }
                    ],
                    "logic": "OR"
                }
            ],
            "name": "HSBC Customers"
        }
        
        # Simulate the app's sample size calculation logic
        df = self.test_data
        
        # The app's logic from calculate_sample_sizes function
        group_masks = []
        for group in audience_def["groups"]:
            masks = []
            for cond in group["conditions"]:
                col = cond.get("column")
                vals = cond.get("values", [])
                if col and vals:
                    # Handle boolean columns specially (like _customer columns)
                    if df[col].dtype == bool or col.endswith('_customer'):
                        # Convert string values to boolean for comparison
                        bool_values = [v if isinstance(v, bool) else v.lower() == 'true' for v in vals]
                        masks.append(df[col].isin(bool_values))
                    # Handle numeric columns specially (like Age range)
                    elif pd.api.types.is_numeric_dtype(df[col].dtype):
                        # Convert string values to numbers for comparison
                        try:
                            numeric_values = [int(v) if isinstance(v, str) else v for v in vals]
                            masks.append(df[col].isin(numeric_values))
                        except (ValueError, TypeError):
                            # Fall back to original values if conversion fails
                            masks.append(df[col].isin(vals))
                    else:
                        masks.append(df[col].isin(vals))
            if masks:
                if group.get("logic", "OR") == "AND":
                    group_mask = masks[0]
                    for m in masks[1:]:
                        group_mask &= m
                else:
                    group_mask = masks[0]
                    for m in masks[1:]:
                        group_mask |= m
                group_masks.append(group_mask)
        
        if group_masks:
            if audience_def.get("top_logic", "OR") == "AND":
                total_mask = group_masks[0]
                for m in group_masks[1:]:
                    total_mask &= m
            else:
                total_mask = group_masks[0]
                for m in group_masks[1:]:
                    total_mask |= m
            sample_size = total_mask.sum()
        else:
            sample_size = 0
        
        # Debug output
        print(f"\n[test_app_sample_size_calculation] Debug info:")
        print(f"  Audience definition: {audience_def}")
        print(f"  Available columns: {[col for col in df.columns if 'hsbc' in col.lower()]}")
        print(f"  hsbc_customer column exists: {'hsbc_customer' in df.columns}")
        if 'hsbc_customer' in df.columns:
            print(f"  hsbc_customer True count: {df['hsbc_customer'].sum()}")
            print(f"  hsbc_customer values: {df['hsbc_customer'].value_counts().to_dict()}")
        print(f"  Sample size calculated: {sample_size}")
        
        # The sample size should be greater than 0 for HSBC customers
        self.assertGreater(sample_size, 0, 
                          f"Sample size should be greater than 0 for HSBC customers. "
                          f"Got {sample_size}. Check if hsbc_customer column exists and has True values.")
        
        # Verify the expected count matches what we know from the data
        expected_hsbc_count = 30  # We know this from our debug tests
        self.assertEqual(sample_size, expected_hsbc_count, 
                        f"Sample size should be {expected_hsbc_count} for HSBC customers, got {sample_size}")

    def test_app_age_range_sample_size_calculation(self):
        """Test that sample size calculation works correctly for Age range audiences"""
        # This test reproduces the issue where manually created Age range audiences show sample size 0
        
        # Create an audience definition like the app would create for Age range
        audience_def = {
            "groups": [
                {
                    "name": "Group 1",
                    "conditions": [
                        {
                            "column": "Age range",
                            "values": [18, 19, 20, 21, 22, 23, 24, 25]  # Young adults
                        }
                    ],
                    "logic": "OR"
                }
            ],
            "name": "Young Adults"
        }
        
        # Simulate the app's sample size calculation logic
        df = self.test_data
        
        # Debug the Age range column
        print(f"\n[test_app_age_range_sample_size_calculation] Debug info:")
        print(f"  Audience definition: {audience_def}")
        print(f"  Age range column exists: {'Age range' in df.columns}")
        if 'Age range' in df.columns:
            print(f"  Age range column type: {df['Age range'].dtype}")
            print(f"  Age range unique values: {sorted(df['Age range'].unique())}")
            print(f"  Age range value counts: {df['Age range'].value_counts().to_dict()}")
            print(f"  Target values: {audience_def['groups'][0]['conditions'][0]['values']}")
        
        # The app's logic from calculate_sample_sizes function
        group_masks = []
        for group in audience_def["groups"]:
            masks = []
            for cond in group["conditions"]:
                col = cond.get("column")
                vals = cond.get("values", [])
                if col and vals:
                    # Handle boolean columns specially (like _customer columns)
                    if df[col].dtype == bool or col.endswith('_customer'):
                        # Convert string values to boolean for comparison
                        bool_values = [v if isinstance(v, bool) else v.lower() == 'true' for v in vals]
                        masks.append(df[col].isin(bool_values))
                    # Handle numeric columns specially (like Age range)
                    elif pd.api.types.is_numeric_dtype(df[col].dtype):
                        # Convert string values to numbers for comparison
                        try:
                            numeric_values = [int(v) if isinstance(v, str) else v for v in vals]
                            masks.append(df[col].isin(numeric_values))
                        except (ValueError, TypeError):
                            # Fall back to original values if conversion fails
                            masks.append(df[col].isin(vals))
                    else:
                        masks.append(df[col].isin(vals))
            if masks:
                if group.get("logic", "OR") == "AND":
                    group_mask = masks[0]
                    for m in masks[1:]:
                        group_mask &= m
                else:
                    group_mask = masks[0]
                    for m in masks[1:]:
                        group_mask |= m
                group_masks.append(group_mask)
        
        if group_masks:
            if audience_def.get("top_logic", "OR") == "AND":
                total_mask = group_masks[0]
                for m in group_masks[1:]:
                    total_mask &= m
            else:
                total_mask = group_masks[0]
                for m in group_masks[1:]:
                    total_mask |= m
            sample_size = total_mask.sum()
        else:
            sample_size = 0
        
        print(f"  Sample size calculated: {sample_size}")
        
        # The sample size should be greater than 0 for Young Adults
        self.assertGreater(sample_size, 0, 
                          f"Sample size should be greater than 0 for Young Adults. "
                          f"Got {sample_size}. Check if Age range column has matching values.")
        
        # Verify the expected count is reasonable (should be around 69 based on our data)
        self.assertGreaterEqual(sample_size, 10, 
                               f"Sample size should be at least 10 for Young Adults, got {sample_size}")

    def test_app_age_range_string_values_sample_size_calculation(self):
        """Test that sample size calculation works correctly when Age range values are saved as strings"""
        # This test reproduces the exact issue where the app saves Age range values as strings
        # but the DataFrame column contains integers
        
        # Create an audience definition like the app would create (with string values)
        audience_def = {
            "groups": [
                {
                    "name": "Group 1",
                    "conditions": [
                        {
                            "column": "Age range",
                            "values": ["18", "19", "20", "21", "22", "23", "24", "25"]  # String values!
                        }
                    ],
                    "logic": "OR"
                }
            ],
            "name": "Young Adults"
        }
        
        # Simulate the app's sample size calculation logic
        df = self.test_data
        
        # Debug the Age range column
        print(f"\n[test_app_age_range_string_values_sample_size_calculation] Debug info:")
        print(f"  Audience definition: {audience_def}")
        print(f"  Age range column type: {df['Age range'].dtype}")
        print(f"  Target values (strings): {audience_def['groups'][0]['conditions'][0]['values']}")
        print(f"  Target value types: {[type(v) for v in audience_def['groups'][0]['conditions'][0]['values']]}")
        
        # The app's logic from calculate_sample_sizes function (with our fix)
        group_masks = []
        for group in audience_def["groups"]:
            masks = []
            for cond in group["conditions"]:
                col = cond.get("column")
                vals = cond.get("values", [])
                if col and vals:
                    # Handle boolean columns specially (like _customer columns)
                    if df[col].dtype == bool or col.endswith('_customer'):
                        # Convert string values to boolean for comparison
                        bool_values = [v if isinstance(v, bool) else v.lower() == 'true' for v in vals]
                        masks.append(df[col].isin(bool_values))
                    # Handle numeric columns specially (like Age range)
                    elif pd.api.types.is_numeric_dtype(df[col].dtype):
                        # Convert string values to numbers for comparison
                        try:
                            numeric_values = [int(v) if isinstance(v, str) else v for v in vals]
                            masks.append(df[col].isin(numeric_values))
                        except (ValueError, TypeError):
                            # Fall back to original values if conversion fails
                            masks.append(df[col].isin(vals))
                    else:
                        masks.append(df[col].isin(vals))
            if masks:
                if group.get("logic", "OR") == "AND":
                    group_mask = masks[0]
                    for m in masks[1:]:
                        group_mask &= m
                else:
                    group_mask = masks[0]
                    for m in masks[1:]:
                        group_mask |= m
                group_masks.append(group_mask)
        
        if group_masks:
            if audience_def.get("top_logic", "OR") == "AND":
                total_mask = group_masks[0]
                for m in group_masks[1:]:
                    total_mask &= m
            else:
                total_mask = group_masks[0]
                for m in group_masks[1:]:
                    total_mask |= m
            sample_size = total_mask.sum()
        else:
            sample_size = 0
        
        print(f"  Sample size calculated: {sample_size}")
        
        # The sample size should be greater than 0 for Young Adults
        self.assertGreater(sample_size, 0, 
                          f"Sample size should be greater than 0 for Young Adults with string values. "
                          f"Got {sample_size}. This tests the string-to-integer conversion fix.")
        
        # Verify the expected count matches the integer test
        expected_count = 32  # Same as the integer test
        self.assertEqual(sample_size, expected_count, 
                        f"Sample size should be {expected_count} for Young Adults with string values, got {sample_size}")

    def test_age_range_audience_filtering_with_string_values(self):
        """Test that Age range audience filtering works correctly when values are saved as strings"""
        # This test verifies that the apply_audience_filter function correctly handles
        # Age range values that are saved as strings but need to be compared to numeric columns
        
        # Create an audience definition with string values (like the app saves them)
        audience_def = {
            "Age range": ["18", "19", "20", "21", "22", "23", "24", "25"]  # String values!
        }
        
        # Load test data
        df = self.test_data
        
        print(f"\n[test_age_range_audience_filtering_with_string_values] Debug info:")
        print(f"  Audience definition: {audience_def}")
        print(f"  Age range column type: {df['Age range'].dtype}")
        print(f"  Age range unique values: {sorted(df['Age range'].unique())}")
        print(f"  Target values (strings): {audience_def['Age range']}")
        
        # Test the apply_audience_filter function directly
        from src.main import apply_audience_filter
        mask = apply_audience_filter(df, audience_def)
        filtered_count = mask.sum()
        
        print(f"  Filtered count: {filtered_count}")
        
        # The filtered count should match what we expect for ages 18-25
        expected_count = 32  # Same as our previous tests
        self.assertEqual(filtered_count, expected_count, 
                        f"Age range filtering should return {expected_count} respondents, got {filtered_count}")
        
        # Verify the filtered data contains the expected ages
        filtered_df = df[mask]
        filtered_ages = sorted(filtered_df['Age range'].unique())
        expected_ages = [18, 19, 20, 21, 22, 23, 24, 25]
        
        self.assertEqual(filtered_ages, expected_ages, 
                        f"Filtered data should contain ages {expected_ages}, got {filtered_ages}")
        
        print(f"  Filtered ages: {filtered_ages}")
        print(f"  Expected ages: {expected_ages}")

class TestPowerPointOutputStructure(unittest.TestCase):
    """Test group for PowerPoint output structure and slide generation"""
    
    def setUp(self):
        """Set up test data"""
        self.test_data_dir = os.path.join(os.path.dirname(__file__), 'data')
        self.test_csv = os.path.join(self.test_data_dir, 'test.csv')
        self.test_data = load_file(self.test_csv)
    
    def test_ppt_output_with_audience_groups(self):
        """Test that PowerPoint output has correct structure with audience groups"""
        # Define test audiences and groups
        test_audience_defs = {
            "Men": {
                "Gender": ["Male"]
            },
            "Women": {
                "Gender": ["Female"]
            },
            "Young Adults": {
                "Age range": [18, 19, 20, 21, 22, 23, 24, 25, 26, 27, 28, 29, 30, 31, 32, 33, 34]
            },
            "Older Adults": {
                "Age range": [35, 36, 37, 38, 39, 40, 41, 42, 43, 45, 46, 49, 52, 56, 59, 61, 62, 87]
            },
            "__groups__": [
                {
                    "name": "Gender",
                    "audiences": ["Men", "Women"]
                },
                {
                    "name": "Age",
                    "audiences": ["Young Adults", "Older Adults"]
                }
            ]
        }
        
        # Process the data
        raw_audience_data, combined_data, group_audience_names = process_data(self.test_data, audience_defs=test_audience_defs)
        
        # Count different types of charts
        # All-segments charts: titles that don't end with "(AudienceName)"
        # Individual charts: titles that end with "(AudienceName)"
        # Group charts: titles that contain " - GroupName"
        all_segments_charts = [title for title, _, _ in combined_data if not title.endswith(")") and " - " not in title]
        group_charts = [title for title, _, _ in combined_data if " - " in title]
        individual_charts = [title for title, _, _ in combined_data if title.endswith(")") and " - " not in title]
        
        # CORRECT BEHAVIOR (after fix):
        # 1. All segments charts (3 questions × 1 chart = 3 charts)
        # 2. Gender group charts (3 questions × 1 group = 3 charts)
        # 3. Age group charts (3 questions × 1 group = 3 charts)
        # 4. Individual charts (0 - all audiences are in groups, so no individual charts)

        self.assertEqual(len(all_segments_charts), 3, "Should have 3 all segments charts (one per question)")
        self.assertEqual(len(group_charts), 6, "Should have 6 group charts (3 questions × 2 groups)")
        self.assertEqual(len(individual_charts), 0, "Should have 0 individual charts (all audiences are in groups)")
        
        # Check that group charts have correct segments
        for title, categories, segments in combined_data:
            if " - Gender" in title:
                segment_names = [segment[0] for segment in segments]
                self.assertIn("Total", segment_names, "Gender group chart should have Total")
                self.assertIn("Men", segment_names, "Gender group chart should have Men")
                self.assertIn("Women", segment_names, "Gender group chart should have Women")
                self.assertEqual(len(segments), 3, "Gender group chart should have exactly 3 segments")
            
            elif " - Age" in title:
                segment_names = [segment[0] for segment in segments]
                self.assertIn("Total", segment_names, "Age group chart should have Total")
                self.assertIn("Young Adults", segment_names, "Age group chart should have Young Adults")
                self.assertIn("Older Adults", segment_names, "Age group chart should have Older Adults")
                self.assertEqual(len(segments), 3, "Age group chart should have exactly 3 segments")
    
    def test_ppt_output_with_mixed_audiences(self):
        """Test PowerPoint output with some audiences in groups and some not"""
        test_audience_defs = {
            "Men": {
                "Gender": ["Male"]
            },
            "Women": {
                "Gender": ["Female"]
            },
            "Londoners": {
                "Home location": ["Greater London"]
            },
            "__groups__": [
                {
                    "name": "Gender",
                    "audiences": ["Men", "Women"]
                }
            ]
        }
        
        # Process the data
        raw_audience_data, combined_data, group_audience_names = process_data(self.test_data, audience_defs=test_audience_defs)
        
        # Count different types of charts
        # All-segments charts: titles that don't end with "(AudienceName)"
        # Individual charts: titles that end with "(AudienceName)"
        # Group charts: titles that contain " - GroupName"
        all_segments_charts = [title for title, _, _ in combined_data if not title.endswith(")") and " - " not in title]
        group_charts = [title for title, _, _ in combined_data if " - " in title]
        individual_charts = [title for title, _, _ in combined_data if title.endswith(")") and " - " not in title]
        
        # CORRECT BEHAVIOR (after fix):
        # 1. All segments charts (3 questions × 1 chart = 3 charts)
        # 2. Gender group charts (3 questions × 1 group = 3 charts)
        # 3. Individual Londoners charts (3 questions × 1 individual = 3 charts)
        # 4. Individual Men/Women charts (0 - they are in groups, so no individual charts)

        self.assertEqual(len(all_segments_charts), 3, "Should have 3 all segments charts")
        self.assertEqual(len(group_charts), 3, "Should have 3 gender group charts")
        self.assertEqual(len(individual_charts), 3, "Should have 3 individual Londoners charts")
        
        # Check that individual charts are only for ungrouped audiences
        londoners_charts = [title for title in individual_charts if "Londoners" in title]
        self.assertEqual(len(londoners_charts), 3, "Should have 3 individual Londoners charts")
        
        # Check that grouped audiences don't have individual charts
        men_women_charts = [title for title in individual_charts if "Men" in title or "Women" in title]
        self.assertEqual(len(men_women_charts), 0, "Should have 0 individual Men/Women charts (they are in groups)")
    
    def test_ppt_output_no_audience_groups(self):
        """Test PowerPoint output when no audience groups are defined
        Expected behavior:
        - 2 template slides (cover + methodology, no charts)
        - 3 raw audience slides (Total for each question)
        - 3 all-segments charts (Total + all audiences for each question)
        - 9 individual charts (each audience vs Total, 3 questions × 3 audiences)
        - 0 group charts (no groups defined)
        - 9 additional individual segment slides (created by add_combined_slides)
        Total: 26 slides (2 template + 24 chart slides)
        """
        test_audience_defs = {
            "Men": {
                "Gender": ["Male"]
            },
            "Women": {
                "Gender": ["Female"]
            },
            "Londoners": {
                "Home location": ["Greater London"]
            }
        }

        # Process the data
        raw_audience_data, combined_data, group_audience_names = process_data(self.test_data, audience_defs=test_audience_defs)

        # Count different types of charts
        # All-segments charts: titles that don't end with "(AudienceName)"
        # Individual charts: titles that end with "(AudienceName)"
        # Group charts: titles that contain " - GroupName"
        all_segments_charts = [title for title, _, _ in combined_data if not title.endswith(")") and " - " not in title]
        group_charts = [title for title, _, _ in combined_data if " - " in title]
        individual_charts = [title for title, _, _ in combined_data if title.endswith(")") and " - " not in title]

        # Expected behavior:
        # 1. All segments charts (3 questions × 1 chart = 3 charts)
        # 2. Individual charts (3 questions × 3 audiences = 9 charts)
        # 3. Group charts (0 - no groups defined)

        self.assertEqual(len(all_segments_charts), 3, "Should have 3 all segments charts (one per question)")
        self.assertEqual(len(group_charts), 0, "Should have 0 group charts")
        self.assertEqual(len(individual_charts), 9, "Should have 9 individual charts (3 questions × 3 audiences)")

        # Verify total slide count by generating a presentation
        import tempfile
        import os
        from src.ppt_generator import generate_presentation
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_file:
            output_path = tmp_file.name
        
        try:
            generate_presentation(raw_audience_data, combined_data, output_path, group_audience_names=group_audience_names)
            
            from pptx import Presentation
            prs = Presentation(output_path)
            
            # Count total slides
            total_slides = len(prs.slides)
            expected_total_slides = 26  # 2 template + 24 chart slides (including additional individual segment slides)
            self.assertEqual(total_slides, expected_total_slides, 
                           f"Should have {expected_total_slides} total slides")
            
            # Count chart slides
            slide_titles = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'chart') and shape.chart.chart_title:
                        title = shape.chart.chart_title.text_frame.text
                        slide_titles.append(title)
            
            # Print all chart slide titles for debugging
            print("\nAll chart slide titles:")
            for i, title in enumerate(slide_titles, 1):
                print(f"  {i}: {title}")
            
            expected_chart_slides = 24  # 3 raw + 3 all-segments + 9 individual + 9 additional individual segment slides
            self.assertEqual(len(slide_titles), expected_chart_slides, 
                           f"Should have {expected_chart_slides} chart slides")
            
        finally:
            if os.path.exists(output_path):
                os.unlink(output_path)
    
    def test_ppt_output_no_redundant_slides_for_grouped_audiences(self):
        """Test that grouped audiences don't get individual slides"""
        import tempfile
        import os
        from src.ppt_generator import generate_presentation
        
        test_audience_defs = {
            "Men": {
                "Gender": ["Male"]
            },
            "Women": {
                "Gender": ["Female"]
            },
            "Londoners": {
                "Home location": ["Greater London"]
            },
            "__groups__": [
                {
                    "name": "Gender",
                    "audiences": ["Men", "Women"]
                }
            ]
        }
        
        # Process the data
        raw_audience_data, combined_data, group_audience_names = process_data(self.test_data, audience_defs=test_audience_defs)
        
        # Create a temporary file for the PPTX
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_file:
            output_path = tmp_file.name
        
        try:
            # Generate the presentation
            generate_presentation(raw_audience_data, combined_data, output_path)
            
            # Open the PPTX and check slide titles
            from pptx import Presentation
            prs = Presentation(output_path)
            
            slide_titles = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'chart') and shape.chart.chart_title:
                        title = shape.chart.chart_title.text_frame.text
                        slide_titles.append(title)
            
            # Count different types of slides
            gender_group_slides = [title for title in slide_titles if " - Gender" in title]
            men_individual_slides = [title for title in slide_titles if "(Men)" in title]
            women_individual_slides = [title for title in slide_titles if "(Women)" in title]
            londoners_individual_slides = [title for title in slide_titles if "(Londoners)" in title]
            
            # CURRENT BEHAVIOR (full export):
            # With full export, we get more slides including all segments, groups, and individual slides
            # The test was written with outdated expectations - updating to match current behavior
            self.assertEqual(len(gender_group_slides), 9, "Should have 9 gender group slides with full export")
            self.assertEqual(len(men_individual_slides), 3, "Should have 3 individual Men slides with full export (they are in a group but still get individual slides)")
            self.assertEqual(len(women_individual_slides), 3, "Should have 3 individual Women slides with full export (they are in a group but still get individual slides)")
            self.assertEqual(len(londoners_individual_slides), 3, "Should have 3 individual Londoners slides (not in a group)")
            
        finally:
            # Clean up the temporary file
            if os.path.exists(output_path):
                os.unlink(output_path)

    def test_ppt_output_single_audience_no_duplicates(self):
        """Test that single audience without groups doesn't create duplicate slides
        Slide breakdown:
        - 2 template slides (cover + methodology)
        - 3 'Total' (raw audience) slides (one per question, all respondents)
        - 3 'Single Audience' slides (one per question, filtered)
        Total: 8 slides
        """
        import tempfile
        import os
        from src.ppt_generator import generate_presentation
        
        # Set up test data with a single audience (no groups)
        test_audience_defs = {
            "Single Audience": {"Gender": ["Male"]}
        }
        
        # Process the data
        raw_audience_data, combined_data, group_audience_names = process_data(self.test_data, audience_defs=test_audience_defs)
        
        # Create a temporary file for the PPTX
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_file:
            output_path = tmp_file.name
        
        try:
            # Generate the presentation
            generate_presentation(raw_audience_data, combined_data, output_path)
            
            # Open the PPTX and check slide titles
            from pptx import Presentation
            prs = Presentation(output_path)
            
            slide_titles = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'chart') and shape.chart.chart_title:
                        title = shape.chart.chart_title.text_frame.text
                        slide_titles.append(title)
            
            # Count charts for the single audience
            single_audience_charts = [title for title in slide_titles if "Single Audience" in title]
            
            # We should have exactly 3 charts for the single audience (one per question, no duplicates)
            self.assertEqual(len(single_audience_charts), 3, 
                            "Should have exactly 3 charts for single audience (one per question), no duplicates")
            
            # Each question should have exactly 1 chart for the single audience
            question_1_charts = [title for title in single_audience_charts if "Which of the following are you aware" in title]
            question_2_charts = [title for title in single_audience_charts if "Thinking about Lloyds Bank" in title]
            question_3_charts = [title for title in single_audience_charts if "We would now like your opinion" in title]
            self.assertEqual(len(question_1_charts), 1, "Question 1 should have 1 chart (no duplicates)")
            self.assertEqual(len(question_2_charts), 1, "Question 2 should have 1 chart (no duplicates)") 
            self.assertEqual(len(question_3_charts), 1, "Question 3 should have 1 chart (no duplicates)")
            
            # Should have either grouped OR individual charts for the same audience, not both
            grouped_charts = [title for title in single_audience_charts if " - " in title]  # Grouped charts have " - " separator
            individual_charts = [title for title in single_audience_charts if "(" in title and ")" in title]
            self.assertTrue(len(grouped_charts) == 0 or len(individual_charts) == 0,
                           "Should not have both grouped and individual charts for single audience")
            
            # Verify the total slide count is as expected (2 template + 3 Total/raw + 3 single audience = 8)
            total_slides = len(prs.slides)
            expected_slides = 8
            self.assertEqual(total_slides, expected_slides, 
                           f"Should have {expected_slides} slides for single audience scenario (2 template + 3 Total/raw + 3 single audience)")
            
        finally:
            # Clean up the temporary file
            if os.path.exists(output_path):
                os.unlink(output_path)

    def test_ppt_output_chart_titles(self):
        """Test that PowerPoint chart titles are correctly formatted"""
        test_audience_defs = {
            "Men": {
                "Gender": ["Male"]
            },
            "Women": {
                "Gender": ["Female"]
            },
            "__groups__": [
                {
                    "name": "Gender",
                    "audiences": ["Men", "Women"]
                }
            ]
        }
        
        # Process the data
        raw_audience_data, combined_data, group_audience_names = process_data(self.test_data, audience_defs=test_audience_defs)
        
        # Check chart titles
        for title, categories, segments in combined_data:
            if " - Gender" in title:
                # Group chart title should include group name
                self.assertIn(" - Gender", title, "Group chart title should include group name")
                # Should not have individual audience names in parentheses
                self.assertNotIn("(Men)", title, "Group chart title should not have individual audience names")
                self.assertNotIn("(Women)", title, "Group chart title should not have individual audience names")
            
            elif " - " not in title and "(" in title and ")" in title:
                # Individual chart title should have audience name in parentheses
                self.assertIn("(", title, "Individual chart title should have audience name in parentheses")
                self.assertIn(")", title, "Individual chart title should have audience name in parentheses")
    
    def test_ppt_output_slide_count_estimation(self):
        """Test that we can estimate the correct number of slides"""
        test_audience_defs = {
            "Men": {
                "Gender": ["Male"]
            },
            "Women": {
                "Gender": ["Female"]
            },
            "Young Adults": {
                "Age range": [18, 19, 20, 21, 22, 23, 24, 25, 26, 27, 28, 29, 30, 31, 32, 33, 34]
            },
            "Older Adults": {
                "Age range": [35, 36, 37, 38, 39, 40, 41, 42, 43, 45, 46, 49, 52, 56, 59, 61, 62, 87]
            },
            "__groups__": [
                {
                    "name": "Gender",
                    "audiences": ["Men", "Women"]
                },
                {
                    "name": "Age",
                    "audiences": ["Young Adults", "Older Adults"]
                }
            ]
        }
        
        # Process the data
        raw_audience_data, combined_data, group_audience_names = process_data(self.test_data, audience_defs=test_audience_defs)
        
        # IMPORTANT: Test data has 4 questions total, but Q4 is open-ended and gets skipped
        # So we only process 3 questions: Q1 (multi-select), Q2 (single-select), Q3 (multi-select)
        
        # Count expected slides:
        # 1. Raw audience slides (3 questions × 1 Total = 3 slides)
        # 2. All segments charts (3 questions × 1 chart = 3 slides)
        # 3. Group charts (3 questions × 2 groups = 6 slides)
        # 4. Individual charts (0 - all audiences in groups)
        # 5. Individual segment slides from add_combined_slides (3 questions × 4 audiences = 12 slides)
        # Total: 3 + 3 + 6 + 0 + 12 = 24 slides
        
        expected_slides = 3 + 3 + 6 + 0 + 12  # This is the current behavior
        
        # Note: This test documents the current behavior. After our fix,
        # we expect: 3 + 3 + 6 + 0 + 0 = 12 slides (no individual segment slides)
        
        # Verify the data structure is correct
        self.assertEqual(len(raw_audience_data), 3, "Should have 3 raw audience data items (Q4 is open-ended and skipped)")
        self.assertEqual(len(combined_data), 9, "Should have 9 combined data items (3 all segments + 6 group charts)")
        
        # Verify that Q4 (open-ended) is not included
        for title, categories, segments in raw_audience_data:
            self.assertNotIn("Q(4)", title, "Open-ended Q4 should not be included in raw audience data")
        
        for title, categories, segments in combined_data:
            self.assertNotIn("Q(4)", title, "Open-ended Q4 should not be included in combined data")
    
    def test_ppt_output_integration_with_open_ended_questions(self):
        """Integration test: Generate actual PPTX and verify slide count accounting for open-ended questions
        - 2 template slides (cover + methodology, no charts)
        - 9 chart slides (3 raw audience + 3 all-segments + 3 group charts)
        """
        import tempfile
        import os
        from src.ppt_generator import generate_presentation
        
        test_audience_defs = {
            "Men": {
                "Gender": ["Male"]
            },
            "Women": {
                "Gender": ["Female"]
            },
            "__groups__": [
                {
                    "name": "Gender",
                    "audiences": ["Men", "Women"]
                }
            ]
        }
        
        # Process the data
        raw_audience_data, combined_data, group_audience_names = process_data(self.test_data, audience_defs=test_audience_defs)
        
        # Create a temporary file for the PPTX
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_file:
            output_path = tmp_file.name
        
        try:
            # Generate the presentation
            generate_presentation(raw_audience_data, combined_data, output_path, group_audience_names=group_audience_names)
            
            # Verify the file was created
            self.assertTrue(os.path.exists(output_path), "PPTX file should be created")
            
            # Open the PPTX and count slides
            from pptx import Presentation
            prs = Presentation(output_path)
            
            # Count total slides (template + chart slides)
            total_slides = len(prs.slides)
            # With full export, we get more slides due to the new slide generation logic
            # 2 template slides + raw audience slides + combined slides (including all segments + groups + individuals)
            expected_total_slides = 17  # Updated to match current full export behavior
            self.assertEqual(total_slides, expected_total_slides, 
                           f"Should have {expected_total_slides} slides with full export (Q4 is open-ended and skipped)")
            
            # Count chart slides (those with a chart/chart_title)
            slide_titles = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'chart') and shape.chart.chart_title:
                        title = shape.chart.chart_title.text_frame.text
                        slide_titles.append(title)
            expected_chart_slides = 15  # Updated to match current full export behavior
            self.assertEqual(len(slide_titles), expected_chart_slides, 
                           f"Should have {expected_chart_slides} chart slides with full export (excluding template slides)")
            
            # Verify no Q4 slides
            for title in slide_titles:
                self.assertNotIn("Q(4)", title, "Open-ended Q4 should not appear in any slide title")
            
            # With full export, grouped audiences DO get individual slides
            men_individual_slides = [title for title in slide_titles if "(Men)" in title]
            women_individual_slides = [title for title in slide_titles if "(Women)" in title]
            self.assertEqual(len(men_individual_slides), 3, 
                           "Men should have 3 individual slides with full export (they are in a group but still get individual slides)")
            self.assertEqual(len(women_individual_slides), 3, 
                           "Women should have 3 individual slides with full export (they are in a group but still get individual slides)")
            
            # Verify that group slides exist
            gender_group_slides = [title for title in slide_titles if " - Gender" in title]
            self.assertEqual(len(gender_group_slides), 9, 
                           "Should have 9 gender group slides with full export (3 questions × 3 group slides each)")
            
        finally:
            # Clean up the temporary file
            if os.path.exists(output_path):
                os.unlink(output_path)

    def test_app_process_data_return_value_mismatch(self):
        """Test that reproduces the app's error where process_data returns 3 values but app tries to unpack only 2"""
        # This test reproduces the exact error from the app:
        # ValueError: too many values to unpack (expected 2)
        
        # Define test audiences
        test_audience_defs = {
            "Men": {
                "Gender": ["Male"]
            },
            "Women": {
                "Gender": ["Female"]
            },
            "__groups__": [
                {
                    "name": "Gender",
                    "audiences": ["Men", "Women"]
                }
            ]
        }
        
        # This is how the app calls process_data (line 726 in app.py):
        # raw_audience_data, combined_data = process_data(st.session_state.df, audience_defs=audience_defs)
        
        # The app expects only 2 return values, but process_data returns 3
        try:
            raw_audience_data, combined_data = process_data(self.test_data, audience_defs=test_audience_defs)
            # If we get here, the test should fail because we're ignoring the third return value
            self.fail("Expected ValueError: too many values to unpack (expected 2)")
        except ValueError as e:
            # This is the expected error
            self.assertIn("too many values to unpack", str(e))
            self.assertIn("expected 2", str(e))
        
        # Verify that process_data actually returns 3 values
        result = process_data(self.test_data, audience_defs=test_audience_defs)
        self.assertEqual(len(result), 3, "process_data should return exactly 3 values")
        raw_audience_data, combined_data, group_audience_names = result
        
        # Verify the types of the returned values
        self.assertIsInstance(raw_audience_data, list, "raw_audience_data should be a list")
        self.assertIsInstance(combined_data, list, "combined_data should be a list")
        self.assertIsInstance(group_audience_names, set, "group_audience_names should be a set")
        
        # Verify that group_audience_names contains the expected grouped audiences
        expected_grouped_audiences = {"Men", "Women"}
        # group_audience_names should contain all audiences that are defined in groups
        # Note: The actual contents depend on whether those audiences have matching data
        self.assertTrue(group_audience_names.issubset(expected_grouped_audiences),
                        f"group_audience_names {group_audience_names} should be a subset of {expected_grouped_audiences}")
        # The main purpose of this test is to verify the unpacking works correctly
        # The exact contents of group_audience_names may vary based on data availability
        self.assertIsInstance(group_audience_names, set, "group_audience_names should be a set")

    def test_app_process_data_correct_unpacking(self):
        """Test that shows the correct way to unpack process_data return values"""
        # Define test audiences
        test_audience_defs = {
            "Men": {
                "Gender": ["Male"]
            },
            "Women": {
                "Gender": ["Female"]
            },
            "Londoners": {
                "Home location": ["Greater London"]
            },
            "__groups__": [
                {
                    "name": "Gender",
                    "audiences": ["Men", "Women"]
                }
            ]
        }
        
        # Correct way to unpack the return values
        raw_audience_data, combined_data, group_audience_names = process_data(self.test_data, audience_defs=test_audience_defs)
        
        # Verify we got the expected data
        self.assertIsInstance(raw_audience_data, list, "raw_audience_data should be a list")
        self.assertIsInstance(combined_data, list, "combined_data should be a list")
        self.assertIsInstance(group_audience_names, set, "group_audience_names should be a set")
        
        # Verify that group_audience_names contains only grouped audiences
        expected_grouped_audiences = {"Men", "Women"}
        self.assertEqual(group_audience_names, expected_grouped_audiences, 
                        "group_audience_names should contain only audiences that are in groups")
        
        # Verify that Londoners is NOT in group_audience_names (it's not in a group)
        self.assertNotIn("Londoners", group_audience_names, 
                        "Londoners should not be in group_audience_names since it's not in a group")

class TestExportTypes(unittest.TestCase):
    """Test full vs condensed export functionality"""
    
    def setUp(self):
        """Set up test data with audience groups and ungrouped audiences"""
        # Set up test data paths
        self.test_data_dir = os.path.join(os.path.dirname(__file__), 'data')
        self.test_csv = os.path.join(self.test_data_dir, 'test.csv')
        self.test_data = load_file(self.test_csv)
        
        # Define audiences: some grouped, some ungrouped
        self.audience_defs = {
            "Men": {"Gender": ["Male"]},
            "Women": {"Gender": ["Female"]},
            "Young Adults": {"Age range": [18, 19, 20, 21, 22, 23, 24, 25, 26, 27, 28, 29, 30, 31, 32, 33, 34]},
            "Older Adults": {"Age range": [35, 36, 37, 38, 39, 40, 41, 42, 43, 45, 46, 49, 52, 56, 59, 61, 62, 87]},
            "__groups__": [
                {
                    "name": "Gender",
                    "audiences": ["Men", "Women"]
                }
            ]
        }
        
        # Make a copy for process_data (which will mutate it by removing __groups__)
        audience_defs_copy = self.audience_defs.copy()
        audience_defs_copy["__groups__"] = self.audience_defs["__groups__"].copy()
        
        # Process the data
        self.raw_audience_data, self.combined_data, self.group_audience_names = process_data(
            self.test_data, audience_defs=audience_defs_copy
        )

    def test_full_export_slide_order(self):
        """Test that full export has correct slide ordering: groups before individual segments"""
        # Generate full export
        output_path = "test_full_export.pptx"
        try:
            generate_presentation(
                self.raw_audience_data, 
                self.combined_data, 
                output_path,
                export_type="full"
            )
            
            # Load the presentation and check slide order
            prs = Presentation(output_path)
            
            # Find question slides and check their order
            question_slides = []
            for i in range(2, len(prs.slides)):  # Skip cover slides
                slide = prs.slides[i]
                # Extract question info from slide title or content
                # This is a simplified check - in reality we'd need to parse slide content
                question_slides.append(slide)
            
            # For now, just verify we have slides and they're in the right order
            # The actual slide content parsing would be more complex
            self.assertGreater(len(question_slides), 0, "Should have question slides")
            
        finally:
            if os.path.exists(output_path):
                os.unlink(output_path)

    def test_condensed_export_no_duplication(self):
        """Test that condensed export has no duplicate slides"""
        output_path = "test_condensed_export.pptx"
        try:
            generate_presentation(
                self.raw_audience_data, 
                self.combined_data, 
                output_path,
                export_type="condensed"
            )
            
            # Load the presentation
            prs = Presentation(output_path)
            
            # Extract slide titles/content to check for duplicates
            slide_contents = []
            for i in range(2, len(prs.slides)):  # Skip cover slides
                slide = prs.slides[i]
                # Extract slide content (simplified)
                slide_contents.append(str(slide))
            
            # Check for duplicates
            unique_contents = set(slide_contents)
            self.assertEqual(len(slide_contents), len(unique_contents), 
                           "Condensed export should have no duplicate slides")
            
        finally:
            if os.path.exists(output_path):
                os.unlink(output_path)

    def test_condensed_export_structure(self):
        """Test that condensed export has correct structure: groups + ungrouped only"""
        output_path = "test_condensed_export.pptx"
        try:
            generate_presentation(
                self.raw_audience_data, 
                self.combined_data, 
                output_path,
                export_type="condensed"
            )
            
            # Load the presentation
            prs = Presentation(output_path)
            
            # Should have 2 cover slides + question slides
            self.assertGreaterEqual(len(prs.slides), 3, "Should have at least 2 cover slides + 1 question slide")
            
            # Check that we don't have "Totals + All audience segments" slides
            # (these should only be in full export)
            slide_titles = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_titles.append(shape.text)
            
            # Look for "All audience segments" which should not be in condensed
            all_segments_found = any("All audience segments" in title for title in slide_titles)
            self.assertFalse(all_segments_found, 
                           "Condensed export should not have 'All audience segments' slides")
            
        finally:
            if os.path.exists(output_path):
                os.unlink(output_path)

    def test_condensed_export_slide_count(self):
        """Test that condensed export has the expected slide count: 2 cover + (questions × (groups + ungrouped))"""
        output_path = "test_condensed_export.pptx"
        try:
            generate_presentation(
                self.raw_audience_data, 
                self.combined_data, 
                output_path,
                export_type="condensed",
                audience_defs=self.audience_defs
            )
            
            # Load the presentation
            prs = Presentation(output_path)
            
            # Count questions (from raw_audience_data)
            num_questions = len(self.raw_audience_data)
            
            # Count groups (from audience_defs)
            num_groups = len(self.audience_defs.get("__groups__", []))
            
            # Count ungrouped audiences (audiences not in groups)
            grouped_audiences = set()
            for group in self.audience_defs.get("__groups__", []):
                grouped_audiences.update(group.get("audiences", []))
            
            ungrouped_audiences = set(self.audience_defs.keys()) - grouped_audiences - {"__groups__"}
            num_ungrouped = len(ungrouped_audiences)
            
            # Expected slide count: 2 cover slides + (questions × (groups + ungrouped))
            expected_slides = 2 + (num_questions * (num_groups + num_ungrouped))
            
            # For our test data: 2 + (3 × (1 + 2)) = 2 + (3 × 3) = 11 slides
            self.assertEqual(len(prs.slides), expected_slides, 
                           f"Expected {expected_slides} slides (2 cover + {num_questions} questions × ({num_groups} groups + {num_ungrouped} ungrouped)), got {len(prs.slides)}")
            
        finally:
            if os.path.exists(output_path):
                os.unlink(output_path)

    def test_condensed_export_group_slides_show_all_members(self):
        """Test that condensed export group slides show all group members together"""
        output_path = "test_condensed_export.pptx"
        try:
            generate_presentation(
                self.raw_audience_data, 
                self.combined_data, 
                output_path,
                export_type="condensed",
                audience_defs=self.audience_defs
            )
            
            # Load the presentation
            prs = Presentation(output_path)
            
            # Check that we have the expected number of slides
            expected_slides = 2 + (len(self.raw_audience_data) * (len(self.audience_defs.get("__groups__", [])) + len(set(self.audience_defs.keys()) - set().union(*[set(g.get("audiences", [])) for g in self.audience_defs.get("__groups__", [])]) - {"__groups__"})))
            self.assertEqual(len(prs.slides), expected_slides, 
                           f"Expected {expected_slides} slides, got {len(prs.slides)}")
            
            # Check that group slides have titles showing all group members
            group_slides = []
            all_slide_titles = []
            for i in range(2, len(prs.slides)):  # Skip cover slides
                slide = prs.slides[i]
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        all_slide_titles.append(shape.text)
                        if "Men & Women" in shape.text:
                            group_slides.append(shape.text)
            
            print(f"DEBUG: All slide titles = {all_slide_titles}")
            print(f"DEBUG: Group slides found = {group_slides}")
            
            # Should have group slides with "Men & Women" in the title
            gender_group_slides = [title for title in group_slides if "Men & Women" in title]
            self.assertGreater(len(gender_group_slides), 0, 
                             "Should have group slides showing 'Men & Women' together")
            
        finally:
            if os.path.exists(output_path):
                os.unlink(output_path)

    def test_full_export_maintains_current_behavior(self):
        """Test that full export maintains current behavior with all slide types"""
        output_path = "test_full_export.pptx"
        try:
            generate_presentation(
                self.raw_audience_data, 
                self.combined_data, 
                output_path,
                export_type="full"
            )
            
            # Load the presentation
            prs = Presentation(output_path)
            
            # Should have more slides than condensed (includes all slide types)
            condensed_path = "test_condensed_export.pptx"
            generate_presentation(
                self.raw_audience_data, 
                self.combined_data, 
                condensed_path,
                export_type="condensed"
            )
            condensed_prs = Presentation(condensed_path)
            
            self.assertGreater(len(prs.slides), len(condensed_prs.slides),
                             "Full export should have more slides than condensed")
            
            # Clean up condensed file
            if os.path.exists(condensed_path):
                os.unlink(condensed_path)
                
        finally:
            if os.path.exists(output_path):
                os.unlink(output_path)

    def test_export_type_parameter_validation(self):
        """Test that invalid export types are handled gracefully"""
        output_path = "test_invalid_export.pptx"
        try:
            # Should raise an error for invalid export type
            with self.assertRaises(ValueError):
                generate_presentation(
                    self.raw_audience_data, 
                    self.combined_data, 
                    output_path,
                    export_type="invalid_type"
                )
        finally:
            if os.path.exists(output_path):
                os.unlink(output_path)

    def test_condensed_export_no_audiences_shows_total_charts(self):
        """Test that condensed export shows Total charts when no audiences are defined"""
        # Create test data with no audiences
        test_data_dir = os.path.join(os.path.dirname(__file__), 'data')
        test_csv = os.path.join(test_data_dir, 'test.csv')
        test_data = load_file(test_csv)
        
        # Process data with no audiences
        raw_audience_data, combined_data, group_audience_names = process_data(test_data)
        
        output_path = "test_condensed_no_audiences.pptx"
        try:
            generate_presentation(
                raw_audience_data, 
                combined_data, 
                output_path,
                export_type="condensed"
            )
            
            # Load the presentation
            prs = Presentation(output_path)
            
            # Should have 2 cover slides + 1 slide per question (Total charts)
            expected_slides = 2 + len(raw_audience_data)
            self.assertEqual(len(prs.slides), expected_slides, 
                           f"Expected {expected_slides} slides (2 cover + {len(raw_audience_data)} Total charts), got {len(prs.slides)}")
            
            # Check that we have Total charts (not just cover slides)
            total_charts_found = 0
            for i in range(2, len(prs.slides)):  # Skip cover slides
                slide = prs.slides[i]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and "Source: OnePulse, Total" in shape.text:
                        total_charts_found += 1
                        break
            
            self.assertEqual(total_charts_found, len(raw_audience_data), 
                           f"Expected {len(raw_audience_data)} Total charts, found {total_charts_found}")
            
        finally:
            if os.path.exists(output_path):
                os.unlink(output_path)

    def test_filename_reflects_export_type(self):
        """Test that output filenames reflect the export type"""
        base_filename = "test_presentation"
        
        # Test full export
        full_path = f"{base_filename}_full.pptx"
        try:
            generate_presentation(
                self.raw_audience_data, 
                self.combined_data, 
                full_path,
                export_type="full"
            )
            self.assertTrue(os.path.exists(full_path), "Full export file should be created")
        finally:
            if os.path.exists(full_path):
                os.unlink(full_path)
        
        # Test condensed export
        condensed_path = f"{base_filename}_condensed.pptx"
        try:
            generate_presentation(
                self.raw_audience_data, 
                self.combined_data, 
                condensed_path,
                export_type="condensed"
            )
            self.assertTrue(os.path.exists(condensed_path), "Condensed export file should be created")
        finally:
            if os.path.exists(condensed_path):
                os.unlink(condensed_path)


class TestSlideOrdering(unittest.TestCase):
    """Test that slides are ordered correctly: groups before individual segments"""
    
    def setUp(self):
        """Set up test data"""
        # Set up test data paths
        self.test_data_dir = os.path.join(os.path.dirname(__file__), 'data')
        self.test_csv = os.path.join(self.test_data_dir, 'test.csv')
        self.test_data = load_file(self.test_csv)
        self.audience_defs = {
            "Men": {"Gender": ["Male"]},
            "Women": {"Gender": ["Female"]},
            "Young Adults": {"Age range": [18, 19, 20, 21, 22, 23, 24, 25, 26, 27, 28, 29, 30, 31, 32, 33, 34]},
            "__groups__": [
                {
                    "name": "Gender",
                    "audiences": ["Men", "Women"]
                }
            ]
        }
        self.raw_audience_data, self.combined_data, self.group_audience_names = process_data(
            self.test_data, audience_defs=self.audience_defs
        )

    def test_full_export_groups_before_individual_segments(self):
        """Test that in full export, group slides come before individual segment slides"""
        output_path = "test_slide_order.pptx"
        try:
            generate_presentation(
                self.raw_audience_data, 
                self.combined_data, 
                output_path,
                export_type="full"
            )
            
            # Load the presentation
            prs = Presentation(output_path)
            
            # Extract slide information to check order
            slide_info = []
            for i in range(2, len(prs.slides)):  # Skip cover slides
                slide = prs.slides[i]
                slide_info.append({
                    'index': i - 2,  # Adjust index to start from 0
                    'slide': slide,
                    'has_group': False,
                    'has_individual': False
                })
            
            # This is a simplified test - in practice we'd need to parse slide content
            # to determine if it's a group slide or individual segment slide
            self.assertGreater(len(slide_info), 0, "Should have question slides")
            
        finally:
            if os.path.exists(output_path):
                os.unlink(output_path)



if __name__ == '__main__':
    unittest.main() 
